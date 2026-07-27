#!/usr/bin/env python3
"""Collect the six current English-manuscript figures in one editable PPTX.

The deck contains one figure per slide. Existing native-shape PowerPoint
sources are reused without rasterizing them. Figure 5 is reconstructed from
the audited per-query source records so that every trajectory segment, match
line, and query marker remains independently editable. Dense measured point
clouds in Figs. 1 and 3 remain raster image plates because representing tens
of thousands of returns as PowerPoint shapes is impractical; their surrounding
labels, panels, separators, descriptors, and annotations remain editable.
"""

from __future__ import annotations

from copy import deepcopy
from io import BytesIO
from pathlib import Path

import numpy as np
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

import generate_scancontext_style_figures as sc


HERE = Path(__file__).resolve().parent
OUT = HERE / "updown_sc_all_english_figures_editable.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5
HEADER_H = 0.54
MARGIN_X = 0.20
MARGIN_BOTTOM = 0.16

FIGURES = (
    (
        "Fig. 1 | Indoor ceiling contamination and dual-envelope motivation",
        HERE / "scancontext_style_figures_editable.pptx",
        0,
    ),
    (
        "Fig. 2 | End-to-end place-recognition and geometric-verification pipeline",
        HERE / "updown_sc_pipeline_editable.pptx",
        0,
    ),
    (
        "Fig. 3 | Measured scan, conventional SC, and UpDown-SC envelopes",
        HERE / "roof_contamination_editable.pptx",
        0,
    ),
    (
        "Fig. 4 | Descriptor detail retained by the complementary envelopes",
        HERE / "descriptor_detail_editable.pptx",
        0,
    ),
    (
        "Fig. 6 | Lower/upper channel-weight ablation",
        HERE / "weight_ablation_editable.pptx",
        0,
    ),
)


def set_white_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = sc.rgb(sc.WHITE)


def add_header(slide, text: str) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.30), Inches(0.08), Inches(SLIDE_W - 0.60), Inches(0.34)
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = sc.rgb(sc.INK)


def clone_picture(source_shape, target_slide, x, y, w, h):
    picture = target_slide.shapes.add_picture(
        BytesIO(source_shape.image.blob), x, y, width=w, height=h
    )
    picture.crop_left = source_shape.crop_left
    picture.crop_right = source_shape.crop_right
    picture.crop_top = source_shape.crop_top
    picture.crop_bottom = source_shape.crop_bottom
    picture.rotation = source_shape.rotation
    picture.name = source_shape.name
    return picture


def clone_shapes_scaled(source_slide, target_slide, scale: float, dx, dy) -> None:
    """Clone slide shapes while preserving editability and z order."""
    for source_shape in source_slide.shapes:
        x = int(dx + source_shape.left * scale)
        y = int(dy + source_shape.top * scale)
        w = max(1, int(source_shape.width * scale))
        h = max(1, int(source_shape.height * scale))

        if source_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            clone_picture(source_shape, target_slide, x, y, w, h)
            continue

        element = deepcopy(source_shape.element)
        target_slide.shapes._spTree.insert_element_before(element, "p:extLst")
        cloned = target_slide.shapes[-1]
        cloned.left = x
        cloned.top = y
        cloned.width = w
        cloned.height = h


def add_source_figure(prs: Presentation, title: str, source_path: Path, index: int) -> None:
    source = Presentation(source_path)
    source_slide = source.slides[index]
    source_w = source.slide_width
    source_h = source.slide_height

    target = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(target)
    add_header(target, title)

    usable_w = Inches(SLIDE_W - 2 * MARGIN_X)
    usable_h = Inches(SLIDE_H - HEADER_H - MARGIN_BOTTOM)
    scale = min(usable_w / source_w, usable_h / source_h)
    content_w = int(source_w * scale)
    content_h = int(source_h * scale)
    dx = int((prs.slide_width - content_w) / 2)
    dy = int(Inches(HEADER_H) + (usable_h - content_h) / 2)
    clone_shapes_scaled(source_slide, target, scale, dx, dy)


def add_cross(slide, x: float, y: float, size: float = 0.085) -> None:
    for x0, y0, x1, y1 in (
        (x - size, y - size, x + size, y + size),
        (x - size, y + size, x + size, y - size),
    ):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x0),
            Inches(y0),
            Inches(x1),
            Inches(y1),
        )
        line.line.color.rgb = sc.rgb(sc.RED)
        line.line.width = Pt(2.2)


def add_query_circle(slide, x: float, y: float, diameter: float = 0.090) -> None:
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x - diameter / 2),
        Inches(y - diameter / 2),
        Inches(diameter),
        Inches(diameter),
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = sc.rgb(sc.BLUE)
    dot.line.color.rgb = sc.rgb(sc.WHITE)
    dot.line.width = Pt(0.35)


def add_figure5(prs: Presentation) -> None:
    loaded = sc.load_all()
    map_xy = np.loadtxt(sc.IH_MAP)[:, 1:3]
    truth_sc, top1_sc, _ = sc.spatial_arrays(loaded["IH"]["SC"])
    truth_ours, top1_ours, _ = sc.spatial_arrays(loaded["IH"]["UpDown-SC"])
    bounds = sc.common_bounds(
        map_xy,
        np.vstack((truth_sc, truth_ours)),
        np.vstack((top1_sc, top1_ours)),
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(slide)
    add_header(slide, "Fig. 5 | Spatial distribution of all 320 eligible IH Top-1 retrievals")

    panel_y = 1.03
    panel_h = 5.80
    panel_w = 5.78
    panel_xs = (0.62, 6.93)

    def map_to_slide(points, x, y, w, h):
        xmin, xmax, ymin, ymax = bounds
        px = x + (points[:, 0] - xmin) / (xmax - xmin) * w
        py = y + h - (points[:, 1] - ymin) / (ymax - ymin) * h
        return np.column_stack((px, py))

    for x, title, rows in (
        (panel_xs[0], "SC + gravity", loaded["IH"]["SC"]),
        (panel_xs[1], "UpDown-SC + gravity", loaded["IH"]["UpDown-SC"]),
    ):
        sc.add_text(slide, x, 0.57, panel_w, 0.34, title, 14, sc.INK, True, PP_ALIGN.CENTER)
        truth, top1, correct = sc.spatial_arrays(rows)
        map_s = map_to_slide(map_xy, x, panel_y, panel_w, panel_h)
        truth_s = map_to_slide(truth, x, panel_y, panel_w, panel_h)
        top1_s = map_to_slide(top1, x, panel_y, panel_w, panel_h)

        for p0, p1 in zip(map_s[:-1], map_s[1:]):
            sc.add_line(slide, *p0, *p1, sc.LIGHT, 0.70)
        for query, match, ok in zip(truth_s, top1_s, correct):
            if ok:
                add_query_circle(slide, float(query[0]), float(query[1]))
            else:
                line = sc.add_line(slide, *query, *match, sc.RED, 0.45, 78)
                line.line.transparency = 78
                add_cross(slide, float(query[0]), float(query[1]))

        sc.add_text(
            slide,
            x + 0.08,
            panel_y + 0.06,
            1.72,
            0.34,
            f"R@1 = {100.0 * correct.mean():.1f}%",
            13,
            sc.INK,
            True,
        )

    # Shared editable legend.
    add_query_circle(slide, 4.42, 7.22, 0.12)
    sc.add_text(slide, 4.55, 7.09, 1.36, 0.27, "correct query", 10, sc.INK)
    add_cross(slide, 6.31, 7.22, 0.065)
    sc.add_text(slide, 6.43, 7.09, 1.62, 0.27, "wrong Top-1", 10, sc.INK)


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.title = "UpDown-SC English manuscript figures — editable sources"
    prs.core_properties.subject = "One current manuscript figure per slide"
    prs.core_properties.author = "Anonymous authors"

    for title, source, index in FIGURES[:4]:
        add_source_figure(prs, title, source, index)
    add_figure5(prs)
    title, source, index = FIGURES[4]
    add_source_figure(prs, title, source, index)

    prs.save(OUT)
    print(f"Wrote {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
