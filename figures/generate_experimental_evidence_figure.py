#!/usr/bin/env python3
"""Generate quantitative ablation and continuous-localization paper panels.

Figure contract
---------------
Core conclusion:
    Gravity canonicalization is most useful under mixed attitude, the adaptive
    map split removes manual per-scene selection without claiming a uniform
    recall gain, and retrieval initializes a stable continuous localization
    trajectory on the public Construction Hall sequence.
Evidence:
    Every value is read from an archived experiment CSV or trajectory file.
    No simulated observations are used.
Outputs:
    Three label-free manuscript subfigures (PDF/SVG/PNG), a native-shape
    editable PowerPoint, a long-form source-data CSV, and audit metadata.
"""

from __future__ import annotations

import csv
import json
import math
from pathlib import Path

import matplotlib as mpl
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("${UPDOWN_SC_ROOT}/icra2027_runtime/experiments")
OUT = Path(__file__).resolve().parent

ADAPTIVE = ROOT / "adaptive_split_conditional_20260720"
WEIGHTED = ROOT / "updown_weight_ablation_real_20260721"
GRAVITY_INPUTS = {
    ("IH", "Native"): WEIGHTED / "selected/ih_native/summary.csv",
    ("IH", "+G"): WEIGHTED / "selected/ih_gravity/summary.csv",
    ("CH", "Native"): WEIGHTED / "selected/ch_native/summary.csv",
    ("CH", "+G"): WEIGHTED / "selected/ch_gravity/summary.csv",
}
SPLIT_INPUTS = {
    ("IH", "Fixed 2.5 m"): WEIGHTED / "fixed25/ih/summary.csv",
    ("IH", "Adaptive"): WEIGHTED / "selected/ih_gravity/summary.csv",
    ("CH", "Fixed 2.5 m"): WEIGHTED / "fixed25/ch/summary.csv",
    ("CH", "Adaptive"): WEIGHTED / "selected/ch_gravity/summary.csv",
    ("H1→H2", "Fixed 2.5 m"): WEIGHTED / "fixed25/h1_h2/summary.csv",
    ("H1→H2", "Adaptive"): WEIGHTED / "selected/h1_h2/summary.csv",
    ("H1→V1", "Fixed 2.5 m"): WEIGHTED / "fixed25/h1_v1/summary.csv",
    ("H1→V1", "Adaptive"): WEIGHTED / "selected/h1_v1/summary.csv",
}

CH_ROOT = ROOT / "rtk_slam_construction_hall"
CH_2M_ROOT = ROOT / "rtk_slam_construction_hall_2m"
MAP_TRAJECTORY = CH_ROOT / "seq1/session/optimized_poses_tum.txt"
LOCALIZATION = CH_ROOT / "localization/seq2_trusted_pose.csv"
REFERENCE = CH_ROOT / "derived/seq2_truth_in_seq1_map_tum.txt"
EVALUATION = CH_ROOT / "localization/trajectory_evaluation.json"

BLUE = "#2B6CB0"
BLUE_LIGHT = "#BFD7EA"
ORANGE = "#D97706"
TEAL = "#168A7A"
PURPLE = "#7353BA"
INK = "#18212B"
MID = "#66717E"
LIGHT = "#E8EDF2"
GRID = "#D5DCE3"
WHITE = "#FFFFFF"

mpl.rcParams.update(
    {
        "font.family": "sans-serif",
        "font.sans-serif": ["Arial", "Helvetica", "DejaVu Sans", "sans-serif"],
        "font.size": 7.2,
        "axes.labelsize": 7.2,
        "axes.titlesize": 8.0,
        "axes.linewidth": 0.75,
        "xtick.labelsize": 6.8,
        "ytick.labelsize": 6.8,
        "legend.fontsize": 6.8,
        "legend.frameon": False,
        "svg.fonttype": "none",
        "pdf.fonttype": 42,
    }
)


def read_first_row(path: Path) -> dict[str, str]:
    with path.open(newline="") as handle:
        return next(csv.DictReader(handle))


def recall(path: Path) -> tuple[float, float]:
    row = read_first_row(path)
    return 100.0 * float(row["recall_at_1"]), 100.0 * float(row["recall_at_5"])


def load_ablation_data():
    gravity = {
        key: recall(path)
        for key, path in GRAVITY_INPUTS.items()
    }
    split = {
        key: recall(path)
        for key, path in SPLIT_INPUTS.items()
    }

    expected_gravity = {
        ("IH", "Native"): (69.0625, 86.875),
        ("IH", "+G"): (69.375, 86.875),
        ("CH", "Native"): (51.351351351, 70.270270270),
        ("CH", "+G"): (64.189189189, 87.162162162),
    }
    for key, expected in expected_gravity.items():
        if not np.allclose(gravity[key], expected, atol=1e-7):
            raise RuntimeError(f"Unexpected gravity result for {key}: {gravity[key]}")

    expected_split = {
        ("IH", "Fixed 2.5 m"): (68.4375, 87.1875),
        ("IH", "Adaptive"): (69.375, 86.875),
        ("CH", "Fixed 2.5 m"): (64.189189189, 86.486486486),
        ("CH", "Adaptive"): (64.189189189, 87.162162162),
        ("H1→H2", "Fixed 2.5 m"): (48.275862069, 68.965517241),
        ("H1→H2", "Adaptive"): (48.275862069, 70.689655172),
        ("H1→V1", "Fixed 2.5 m"): (54.385964912, 80.701754386),
        ("H1→V1", "Adaptive"): (54.385964912, 82.456140351),
    }
    for key, expected in expected_split.items():
        if not np.allclose(split[key], expected, atol=1e-7):
            raise RuntimeError(f"Unexpected split result for {key}: {split[key]}")
    return gravity, split


def load_localization_data():
    map_traj = np.loadtxt(MAP_TRAJECTORY)
    reference = np.loadtxt(REFERENCE)
    estimated = np.genfromtxt(
        LOCALIZATION,
        delimiter=",",
        names=True,
        dtype=None,
        encoding=None,
    )
    stamps = estimated["stamp_ns"].astype(float) / 1e9
    indices = np.searchsorted(stamps, reference[:, 0])
    indices = np.clip(indices, 1, len(stamps) - 1)
    lower = indices - 1
    choose_lower = (
        np.abs(stamps[lower] - reference[:, 0])
        <= np.abs(stamps[indices] - reference[:, 0])
    )
    indices = np.where(choose_lower, lower, indices)
    xy = np.column_stack((estimated["x"], estimated["y"]))
    errors = np.linalg.norm(xy[indices] - reference[:, 1:3], axis=1)
    with EVALUATION.open() as handle:
        summary = json.load(handle)
    checks = (
        (np.median(errors), summary["position_error_median"]),
        (np.quantile(errors, 0.95), summary["position_error_p95"]),
        (np.max(errors), summary["position_error_max"]),
    )
    if any(not math.isclose(a, b, abs_tol=1e-9) for a, b in checks):
        raise RuntimeError("Recomputed localization errors disagree with the audit JSON")
    return map_traj, reference, estimated, errors, summary


def clean_axes(ax):
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.tick_params(length=2.5, width=0.7, color=MID)
    ax.grid(axis="y", color=GRID, linewidth=0.55, alpha=0.85, zorder=0)


def save_panel(fig, stem: str):
    fig.savefig(OUT / f"{stem}.pdf", bbox_inches="tight", pad_inches=0.025)
    fig.savefig(OUT / f"{stem}.svg", bbox_inches="tight", pad_inches=0.025)
    fig.savefig(OUT / f"{stem}.png", dpi=400, bbox_inches="tight", pad_inches=0.025)
    plt.close(fig)


def make_gravity_panel(gravity):
    labels = ["IH R@1", "IH R@5", "CH R@1", "CH R@5"]
    native = [
        gravity[("IH", "Native")][0],
        gravity[("IH", "Native")][1],
        gravity[("CH", "Native")][0],
        gravity[("CH", "Native")][1],
    ]
    aligned = [
        gravity[("IH", "+G")][0],
        gravity[("IH", "+G")][1],
        gravity[("CH", "+G")][0],
        gravity[("CH", "+G")][1],
    ]
    x = np.arange(len(labels))
    width = 0.34
    fig, ax = plt.subplots(figsize=(2.05, 1.78))
    ax.bar(
        x - width / 2,
        native,
        width,
        label="Native",
        facecolor=WHITE,
        edgecolor=MID,
        linewidth=0.9,
        zorder=3,
    )
    ax.bar(
        x + width / 2,
        aligned,
        width,
        label="+G",
        facecolor=BLUE,
        edgecolor=BLUE,
        linewidth=0.8,
        zorder=3,
    )
    for idx, (a, b) in enumerate(zip(native, aligned)):
        delta = b - a
        ax.text(
            idx,
            max(a, b) + 1.2,
            f"{delta:+.1f}",
            ha="center",
            va="bottom",
            fontsize=6.2,
            color=TEAL if delta > 0 else MID,
        )
    ax.set_ylim(45, 94)
    ax.set_ylabel("Recall (%)")
    ax.set_xticks(x, labels, rotation=21, ha="right")
    ax.set_title("Gravity canonicalization")
    ax.legend(loc="lower right", ncol=2, handlelength=1.2, columnspacing=0.8)
    clean_axes(ax)
    fig.tight_layout(pad=0.35)
    save_panel(fig, "experimental_gravity_ablation")


def make_split_panel(split):
    datasets = ["IH", "CH", "H1→H2", "H1→V1"]
    r1_delta = [
        split[(name, "Adaptive")][0] - split[(name, "Fixed 2.5 m")][0]
        for name in datasets
    ]
    r5_delta = [
        split[(name, "Adaptive")][1] - split[(name, "Fixed 2.5 m")][1]
        for name in datasets
    ]
    x = np.arange(len(datasets))
    width = 0.34
    fig, ax = plt.subplots(figsize=(2.08, 1.78))
    ax.axhline(0, color=MID, linewidth=0.75, zorder=1)
    ax.bar(
        x - width / 2,
        r1_delta,
        width,
        color=ORANGE,
        edgecolor=ORANGE,
        label="R@1",
        zorder=3,
    )
    ax.bar(
        x + width / 2,
        r5_delta,
        width,
        color=PURPLE,
        edgecolor=PURPLE,
        label="R@5",
        zorder=3,
    )
    for xpos, value in zip(x - width / 2, r1_delta):
        ax.text(
            xpos,
            value + (0.16 if value >= 0 else -0.16),
            f"{value:+.1f}",
            ha="center",
            va="bottom" if value >= 0 else "top",
            fontsize=5.8,
            color=INK,
        )
    for xpos, value in zip(x + width / 2, r5_delta):
        ax.text(
            xpos,
            value + (0.16 if value >= 0 else -0.16),
            f"{value:+.1f}",
            ha="center",
            va="bottom" if value >= 0 else "top",
            fontsize=5.8,
            color=INK,
        )
    ax.set_ylim(-3.0, 3.8)
    ax.set_ylabel("Adaptive − fixed (points)")
    ax.set_xticks(x, datasets, rotation=18, ha="right")
    ax.set_title("Map-adaptive split")
    ax.legend(loc="lower left", ncol=2, handlelength=1.1, columnspacing=0.8)
    clean_axes(ax)
    fig.tight_layout(pad=0.35)
    save_panel(fig, "experimental_split_ablation")


def make_trajectory_panel(map_traj, reference, estimated, errors, summary):
    fig, ax = plt.subplots(figsize=(3.0, 1.78))
    ax.plot(
        map_traj[:, 1],
        map_traj[:, 2],
        color="#AAB4BF",
        linewidth=1.0,
        label="Map traversal",
        zorder=1,
    )
    ax.plot(
        reference[:, 1],
        reference[:, 2],
        color=ORANGE,
        linewidth=1.3,
        linestyle=(0, (3, 1.8)),
        label="Reference",
        zorder=3,
    )
    ax.plot(
        estimated["x"],
        estimated["y"],
        color=BLUE,
        linewidth=1.0,
        alpha=0.95,
        label="Localized",
        zorder=2,
    )
    ax.scatter(
        reference[0, 1],
        reference[0, 2],
        marker="o",
        s=20,
        facecolor=TEAL,
        edgecolor=WHITE,
        linewidth=0.6,
        zorder=5,
    )
    ax.text(
        0.98,
        0.04,
        (
            f"{len(errors)} reference poses\n"
            f"median {summary['position_error_median']:.3f} m  |  "
            f"95th {summary['position_error_p95']:.3f} m"
        ),
        transform=ax.transAxes,
        ha="right",
        va="bottom",
        fontsize=6.3,
        color=INK,
        bbox=dict(
            boxstyle="round,pad=0.24",
            facecolor=WHITE,
            edgecolor=GRID,
            linewidth=0.6,
            alpha=0.94,
        ),
    )
    ax.set_aspect("equal", adjustable="datalim")
    ax.set_xlabel("x (m)")
    ax.set_ylabel("y (m)")
    ax.set_title("Continuous CH localization")
    ax.legend(loc="upper center", ncol=3, handlelength=1.8, columnspacing=0.8)
    clean_axes(ax)
    fig.tight_layout(pad=0.35)
    save_panel(fig, "experimental_ch_localization")


def export_source_data(gravity, split, map_traj, reference, estimated, errors):
    path = OUT / "experimental_evidence_source_data.csv"
    with path.open("w", newline="") as handle:
        writer = csv.DictWriter(
            handle,
            fieldnames=[
                "panel",
                "dataset",
                "condition",
                "metric",
                "value",
                "x",
                "y",
                "stamp",
            ],
        )
        writer.writeheader()
        for (dataset, condition), values in gravity.items():
            for metric, value in zip(("R@1", "R@5"), values):
                writer.writerow(
                    {
                        "panel": "gravity",
                        "dataset": dataset,
                        "condition": condition,
                        "metric": metric,
                        "value": f"{value:.9f}",
                    }
                )
        for (dataset, condition), values in split.items():
            for metric, value in zip(("R@1", "R@5"), values):
                writer.writerow(
                    {
                        "panel": "split",
                        "dataset": dataset,
                        "condition": condition,
                        "metric": metric,
                        "value": f"{value:.9f}",
                    }
                )
        for label, array, stamp_col, x_col, y_col in (
            ("map", map_traj, 0, 1, 2),
            ("reference", reference, 0, 1, 2),
        ):
            for row in array:
                writer.writerow(
                    {
                        "panel": "trajectory",
                        "condition": label,
                        "x": f"{row[x_col]:.9f}",
                        "y": f"{row[y_col]:.9f}",
                        "stamp": f"{row[stamp_col]:.9f}",
                    }
                )
        stride = max(1, len(estimated) // 1800)
        for index in range(0, len(estimated), stride):
            writer.writerow(
                {
                    "panel": "trajectory",
                    "condition": "localized",
                    "x": f"{float(estimated['x'][index]):.9f}",
                    "y": f"{float(estimated['y'][index]):.9f}",
                    "stamp": f"{float(estimated['stamp_ns'][index]) / 1e9:.9f}",
                }
            )
        for value in errors:
            writer.writerow(
                {
                    "panel": "trajectory_error",
                    "condition": "xy_error",
                    "metric": "error_m",
                    "value": f"{value:.9f}",
                }
            )


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.lstrip("#")
    return RGBColor.from_string(value.upper())


def add_text(slide, x, y, w, h, text, size=8, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Pt(0)
    frame.margin_top = frame.margin_bottom = Pt(0)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_line(slide, x1, y1, x2, y2, color=INK, width=1.0):
    line = slide.shapes.add_connector(
        1,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    return line


def add_rect(slide, x, y, w, h, fill, line=None, line_width=0.7):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    shape.line.width = Pt(line_width)
    return shape


def build_editable_ppt(gravity, split, map_traj, reference, estimated, summary):
    prs = Presentation()
    prs.slide_width = Inches(13.6)
    prs.slide_height = Inches(4.25)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(WHITE)

    panels = [(0.25, 0.25, 3.75, 3.7), (4.15, 0.25, 3.75, 3.7), (8.05, 0.25, 5.3, 3.7)]
    for x, y, w, h in panels:
        add_rect(slide, x, y, w, h, WHITE, GRID, 0.8)

    # Panel 1: native/+G bars.
    add_text(slide, 0.48, 0.42, 3.2, 0.28, "Gravity canonicalization", 11, bold=True)
    categories = ["IH R@1", "IH R@5", "CH R@1", "CH R@5"]
    native = [
        gravity[("IH", "Native")][0],
        gravity[("IH", "Native")][1],
        gravity[("CH", "Native")][0],
        gravity[("CH", "Native")][1],
    ]
    aligned = [
        gravity[("IH", "+G")][0],
        gravity[("IH", "+G")][1],
        gravity[("CH", "+G")][0],
        gravity[("CH", "+G")][1],
    ]
    x0, y0, w0, h0 = 0.72, 0.95, 3.0, 2.25
    add_line(slide, x0, y0 + h0, x0 + w0, y0 + h0, MID, 0.8)
    for i, (label, a, b) in enumerate(zip(categories, native, aligned)):
        cx = x0 + 0.40 + i * 0.72
        for offset, value, fill, outline in (
            (-0.12, a, WHITE, MID),
            (0.12, b, BLUE, BLUE),
        ):
            height = h0 * (value - 45) / 50
            add_rect(slide, cx + offset - 0.095, y0 + h0 - height, 0.19, height, fill, outline)
        add_text(slide, cx - 0.28, y0 + h0 + 0.08, 0.56, 0.35, label, 7, align=PP_ALIGN.CENTER)
        add_text(
            slide,
            cx - 0.25,
            y0 + h0 - h0 * (max(a, b) - 45) / 50 - 0.24,
            0.5,
            0.18,
            f"{b-a:+.1f}",
            7,
            TEAL if b > a else MID,
            True,
            PP_ALIGN.CENTER,
        )
    add_rect(slide, 1.02, 3.64, 0.18, 0.12, WHITE, MID)
    add_text(slide, 1.25, 3.59, 0.55, 0.2, "Native", 7)
    add_rect(slide, 2.15, 3.64, 0.18, 0.12, BLUE, BLUE)
    add_text(slide, 2.38, 3.59, 0.4, 0.2, "+G", 7)

    # Panel 2: adaptive minus fixed split.
    add_text(slide, 4.38, 0.42, 3.2, 0.28, "Map-adaptive split", 11, bold=True)
    names = ["IH", "CH", "H1→H2", "H1→V1"]
    d1 = [split[(name, "Adaptive")][0] - split[(name, "Fixed 2.5 m")][0] for name in names]
    d5 = [split[(name, "Adaptive")][1] - split[(name, "Fixed 2.5 m")][1] for name in names]
    base_y = 2.08
    add_line(slide, 4.55, base_y, 7.65, base_y, MID, 0.8)
    for i, (name, v1, v5) in enumerate(zip(names, d1, d5)):
        cx = 4.95 + i * 0.72
        for offset, value, color in ((-0.12, v1, ORANGE), (0.12, v5, PURPLE)):
            height = abs(value) * 0.35
            top = base_y - height if value >= 0 else base_y
            add_rect(slide, cx + offset - 0.095, top, 0.19, max(height, 0.018), color, color)
            add_text(
                slide,
                cx + offset - 0.22,
                top - 0.22 if value >= 0 else top + height + 0.03,
                0.44,
                0.18,
                f"{value:+.1f}",
                6.5,
                INK,
                align=PP_ALIGN.CENTER,
            )
        add_text(slide, cx - 0.33, 3.05, 0.66, 0.28, name, 7, align=PP_ALIGN.CENTER)
    add_text(slide, 4.48, 0.83, 3.0, 0.22, "Adaptive − fixed 2.5 m (percentage points)", 7, MID)
    add_rect(slide, 5.10, 3.62, 0.18, 0.12, ORANGE)
    add_text(slide, 5.33, 3.57, 0.42, 0.2, "R@1", 7)
    add_rect(slide, 6.08, 3.62, 0.18, 0.12, PURPLE)
    add_text(slide, 6.31, 3.57, 0.42, 0.2, "R@5", 7)

    # Panel 3: native trajectory shapes.
    add_text(slide, 8.28, 0.42, 4.6, 0.28, "Continuous CH localization", 11, bold=True)
    arrays = [
        (map_traj[:, 1:3], "#AAB4BF", 1.0),
        (reference[:, 1:3], ORANGE, 1.5),
        (np.column_stack((estimated["x"], estimated["y"])), BLUE, 1.1),
    ]
    all_xy = np.vstack([array for array, _, _ in arrays])
    minimum = all_xy.min(axis=0)
    maximum = all_xy.max(axis=0)
    span = np.maximum(maximum - minimum, 1e-9)
    px, py, pw, ph = 8.40, 0.88, 4.55, 2.55
    for array, color, width in arrays:
        stride = max(1, len(array) // 850)
        points = array[::stride]
        x = px + (points[:, 0] - minimum[0]) / span[0] * pw
        y = py + ph - (points[:, 1] - minimum[1]) / span[1] * ph
        for i in range(len(points) - 1):
            add_line(slide, x[i], y[i], x[i + 1], y[i + 1], color, width)
    add_text(
        slide,
        9.07,
        3.50,
        3.85,
        0.23,
        (
            f"748 reference poses   median {summary['position_error_median']:.3f} m   "
            f"95th {summary['position_error_p95']:.3f} m"
        ),
        7.2,
        INK,
        True,
        PP_ALIGN.CENTER,
    )
    for x, color, label in ((8.48, "#AAB4BF", "Map"), (9.56, ORANGE, "Reference"), (11.10, BLUE, "Localized")):
        add_line(slide, x, 3.80, x + 0.30, 3.80, color, 2.0)
        add_text(slide, x + 0.36, 3.70, 0.85, 0.22, label, 7)

    prs.save(OUT / "experimental_evidence_editable.pptx")


def write_metadata(summary):
    metadata = {
        "figure_contract": {
            "core_conclusion": (
                "Gravity mainly helps the mixed-attitude CH sequence; the adaptive "
                "split removes manual scene tuning without a uniform recall claim; "
                "the public CH replay remains continuously localized."
            ),
            "archetype": "two quantitative ablations plus trajectory overlay",
            "backend": "Python/Matplotlib with native-shape PowerPoint export",
            "target": "ICRA/IEEE two-column figure*",
        },
        "inputs": {
            "gravity": {f"{k[0]}_{k[1]}": str(v) for k, v in GRAVITY_INPUTS.items()},
            "split": {f"{k[0]}_{k[1]}": str(v) for k, v in SPLIT_INPUTS.items()},
            "map_trajectory": str(MAP_TRAJECTORY),
            "localized_trajectory": str(LOCALIZATION),
            "reference_trajectory": str(REFERENCE),
            "trajectory_audit": str(EVALUATION),
        },
        "trajectory_summary": summary,
        "image_integrity": (
            "Trajectory lines use all input poses in the manuscript raster/vector "
            "panels. The editable PowerPoint uniformly subsamples only the dense "
            "localized line to keep the number of native shapes tractable."
        ),
    }
    with (OUT / "experimental_evidence_metadata.json").open("w") as handle:
        json.dump(metadata, handle, indent=2)


def main():
    gravity, split = load_ablation_data()
    map_traj, reference, estimated, errors, summary = load_localization_data()
    make_gravity_panel(gravity)
    make_split_panel(split)
    make_trajectory_panel(map_traj, reference, estimated, errors, summary)
    export_source_data(gravity, split, map_traj, reference, estimated, errors)
    build_editable_ppt(gravity, split, map_traj, reference, estimated, summary)
    write_metadata(summary)
    print("Generated experimental evidence panels and editable source.")


if __name__ == "__main__":
    main()
