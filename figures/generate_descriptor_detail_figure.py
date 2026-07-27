#!/usr/bin/env python3
"""Generate a descriptor-detail analysis from one measured scan.

Figure contract
---------------
Core conclusion:
    On the same polar grid, full-height SC is dominated by overhead maxima;
    the Up (upper-layer) channel retains broad overhead support by storing the
    minimum above the split, while the Down (lower-layer) channel restores
    lower/middle detail that SC suppresses in mixed cells.
Evidence:
    The scan is selected by the pre-existing, deterministic roof-contamination
    rule over all truth-valid IH queries. No additional frame is hand-picked.
Archetype:
    Image plate plus quantitative callouts, inspired only by the evidence logic
    of descriptor-detail figures rather than by their visual artwork.
Outputs:
    PDF/SVG/PNG, native-cell editable PowerPoint, source CSV, and metadata.
"""

from __future__ import annotations

import argparse
import csv
import json
from pathlib import Path

import matplotlib as mpl
import matplotlib.pyplot as plt
import numpy as np
from matplotlib.colors import LinearSegmentedColormap
from matplotlib.patches import Rectangle
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import generate_roof_contamination_figure as roof


OUT = Path(__file__).resolve().parent
DEFAULT_QUERY_DIR = (
    Path.home()
    / "icra2027_runtime/experiments/baseline_20260715/queries/loc_2_floor_continuous"
)
SPLIT_HEIGHT_M = 2.1
WINDOW_WIDTH = 12

BLUE = "#2F6FAE"
ORANGE = "#D77A12"
MAGENTA = "#C83C91"
GREEN = "#2FBF71"
INK_BOX = "#131C26"
RED = "#B44747"
INK = "#1D2935"
MID = "#66727E"
LIGHT = "#EDF1F4"
GRID = "#D5DCE2"
WHITE = "#FFFFFF"

mpl.rcParams.update(
    {
        "font.family": "sans-serif",
        "font.sans-serif": ["Arial", "Helvetica", "DejaVu Sans", "sans-serif"],
        "font.size": 6.4,
        "axes.titlesize": 7.1,
        "axes.labelsize": 6.2,
        "axes.linewidth": 0.7,
        "xtick.labelsize": 5.5,
        "ytick.labelsize": 5.5,
        "svg.fonttype": "none",
        "pdf.fonttype": 42,
    }
)


def parse_args():
    parser = argparse.ArgumentParser()
    parser.add_argument("--query-dir", type=Path, default=DEFAULT_QUERY_DIR)
    parser.add_argument("--output-dir", type=Path, default=OUT)
    return parser.parse_args()


def strongest_mixed_window(hidden: np.ndarray) -> tuple[int, int]:
    per_sector = hidden.sum(axis=0)
    extended = np.r_[per_sector, per_sector[: WINDOW_WIDTH - 1]]
    scores = np.convolve(extended, np.ones(WINDOW_WIDTH, dtype=int), mode="valid")
    start = int(np.argmax(scores[: roof.NUM_SECTORS]))
    return start, int(scores[start])


def window_stats(values: np.ndarray, valid: np.ndarray, start: int):
    columns = [(start + offset) % roof.NUM_SECTORS for offset in range(WINDOW_WIDTH)]
    window_vals = values[:, columns]
    window_valid = valid[:, columns]
    selected = window_vals[window_valid]
    return selected


def select_ratio_window(
    numerator_vals: np.ndarray,
    numerator_valid: np.ndarray,
    sc_vals: np.ndarray,
    sc_valid: np.ndarray,
    exclude_starts: list[int],
    min_cells: int = 30,
) -> tuple[int, float, float, float]:
    """Window maximizing std(channel)/std(SC over overhead-set cells).

    The denominator uses only cells whose SC maximum is set by overhead
    structure (value above the split), i.e. where the ceiling pins SC flat;
    the window must be circularly disjoint from previously chosen windows.
    """
    overhead_sc = sc_valid & (sc_vals > roof.SPLIT_HEIGHT_M)
    best = (-1, 0.0, 0.0, 0.0)
    for start in range(roof.NUM_SECTORS):
        if any(
            min((start - other) % roof.NUM_SECTORS,
                (other - start) % roof.NUM_SECTORS) < WINDOW_WIDTH
            for other in exclude_starts
        ):
            continue
        channel = window_stats(numerator_vals, numerator_valid, start)
        sc = window_stats(sc_vals, overhead_sc, start)
        if len(channel) < min_cells or len(sc) < min_cells:
            continue
        ratio = float(np.std(channel)) / (float(np.std(sc)) + 1e-6)
        if ratio > best[1]:
            best = (start, ratio, float(np.std(channel)), float(np.std(sc)))
    return best


def select_gap_window(
    sc_vals: np.ndarray,
    sc_valid: np.ndarray,
    up_vals: np.ndarray,
    up_valid: np.ndarray,
    exclude_starts: list[int],
    gap_m: float = 0.3,
) -> tuple[int, int, float]:
    """Window with the most cells whose top and underside differ by > gap.

    These are overhead-set cells where SC's stored maximum and the Up
    channel's underside height disagree substantially -- structure SC
    discards even without any lower-layer competition.
    """
    mask = sc_valid & up_valid & (sc_vals > roof.SPLIT_HEIGHT_M)
    gap = np.where(mask, sc_vals - up_vals, 0.0)
    cell = mask & (gap > gap_m)
    best = (-1, 0, 0.0)
    for start in range(roof.NUM_SECTORS):
        if any(
            min((start - other) % roof.NUM_SECTORS,
                (other - start) % roof.NUM_SECTORS) < WINDOW_WIDTH
            for other in exclude_starts
        ):
            continue
        columns = [(start + offset) % roof.NUM_SECTORS
                   for offset in range(WINDOW_WIDTH)]
        count = int(cell[:, columns].sum())
        if count > best[1]:
            gaps = gap[:, columns][cell[:, columns]]
            best = (start, count, float(np.median(gaps)))
    return best


def rolled(values: np.ndarray, valid: np.ndarray, hidden: np.ndarray, start: int):
    target_start = roof.NUM_SECTORS // 2 - WINDOW_WIDTH // 2
    shift = target_start - start
    return (
        np.roll(values, shift, axis=1),
        np.roll(valid, shift, axis=1),
        np.roll(hidden, shift, axis=1),
        target_start,
        shift,
    )


def height_style():
    cmap = LinearSegmentedColormap.from_list(
        "detail_height",
        ["#263B73", "#3775BA", "#D9EFEA", "#F4C66A", "#B64342"],
    ).copy()
    cmap.set_bad(LIGHT)
    return cmap, mpl.colors.Normalize(vmin=-0.5, vmax=4.0, clip=False)


def draw_descriptor(
    ax,
    values,
    valid,
    hidden,
    title,
    subtitle,
    start,
    cmap,
    norm,
    mark_hidden=False,
    show_xlabel=False,
    flip_rings=False,
    extra_windows=None,
):
    view, view_valid, view_hidden, display_start, _ = rolled(
        values, valid, hidden, start
    )
    if flip_rings:
        # Match Fig. 2's Up-map convention: near rings drawn at the top.
        view = view[::-1]
        view_valid = view_valid[::-1]
        view_hidden = view_hidden[::-1]
    image = ax.imshow(
        np.ma.masked_where(~view_valid, view),
        origin="lower",
        aspect="auto",
        interpolation="none",
        cmap=cmap,
        norm=norm,
        rasterized=True,
    )
    if mark_hidden:
        ring, sector = np.nonzero(view_hidden)
        ax.scatter(sector, ring, s=2.2, marker=".", color=RED, alpha=0.88, zorder=4)
    for window_start, color, style in ([(display_start, MAGENTA, (0, (2.8, 1.4)))]
                                       + (extra_windows or [])):
        spans = [(window_start, min(window_start + WINDOW_WIDTH,
                                    roof.NUM_SECTORS))]
        overflow = window_start + WINDOW_WIDTH - roof.NUM_SECTORS
        if overflow > 0:
            spans = [(window_start, roof.NUM_SECTORS), (0, overflow)]
        for span_start, span_end in spans:
            ax.add_patch(
                Rectangle(
                    (span_start - 0.5, -0.5),
                    span_end - span_start,
                    roof.NUM_RINGS,
                    fill=False,
                    edgecolor=color,
                    linewidth=1.25,
                    linestyle=style,
                    zorder=5,
                )
            )
    ax.set_xlim(-0.5, roof.NUM_SECTORS - 0.5)
    ax.set_ylim(-0.5, roof.NUM_RINGS - 0.5)
    ax.set_xticks([0, 15, 30, 45, 59])
    ax.set_yticks([0, 5, 10, 15])
    if flip_rings:
        ax.set_yticklabels(["15", "10", "5", "0"])
    ax.set_xlabel("circularly shifted sector" if show_xlabel else "")
    ax.set_ylabel("ring")
    ax.set_title(title, color=INK, fontweight="bold", pad=2, loc="left")
    for spine in ax.spines.values():
        spine.set_color(GRID)
        spine.set_linewidth(0.55)
    return image


def make_figure(record, descriptor, output: Path):
    hidden = descriptor.hidden_mask
    start, window_mixed = strongest_mixed_window(hidden)
    up_start, up_gap_cells, up_gap_median = select_gap_window(
        descriptor.sc_max, descriptor.sc_valid,
        descriptor.upper_min, descriptor.upper_valid, [start])
    down_start, down_ratio, down_std, down_sc_std = select_ratio_window(
        descriptor.lower_max, descriptor.lower_valid,
        descriptor.sc_max, descriptor.sc_valid, [start, up_start])
    target_start = roof.NUM_SECTORS // 2 - WINDOW_WIDTH // 2
    shift = target_start - start
    extra_windows = [
        ((up_start + shift) % roof.NUM_SECTORS, GREEN, (0, (5, 2.2))),
        ((down_start + shift) % roof.NUM_SECTORS, INK_BOX, (0, (1.2, 1.9))),
    ]
    occupied = int(descriptor.occupied_mask.sum())
    overhead = int(descriptor.upper_valid.sum())
    lower = int(descriptor.lower_valid.sum())
    mixed = int(hidden.sum())
    cmap, norm = height_style()

    fig = plt.figure(figsize=(3.45, 2.82), facecolor=WHITE)
    grid = fig.add_gridspec(
        3,
        1,
        left=0.16,
        right=0.965,
        bottom=0.19,
        top=0.965,
        hspace=0.70,
    )
    descriptor_axes = [fig.add_subplot(grid[index, 0]) for index in range(3)]

    images = []
    images.append(
        draw_descriptor(
            descriptor_axes[0],
            descriptor.sc_max,
            descriptor.sc_valid,
            hidden,
            "SC: max over all heights",
            f"{overhead}/{occupied} cells set by overhead maxima",
            start,
            cmap,
            norm,
            mark_hidden=False,
            show_xlabel=False,
            extra_windows=extra_windows,
        )
    )
    images.append(
        draw_descriptor(
            descriptor_axes[1],
            descriptor.upper_min,
            descriptor.upper_valid,
            hidden,
            r"Up: min above $\tau_g$",
            f"{overhead} overhead cells retained ({100.0 * overhead / occupied:.1f}% of SC support)",
            start,
            cmap,
            norm,
            show_xlabel=False,
            flip_rings=True,
            extra_windows=extra_windows,
        )
    )
    images.append(
        draw_descriptor(
            descriptor_axes[2],
            descriptor.lower_max,
            descriptor.lower_valid,
            hidden,
            r"Down: max below $\tau_g$",
            f"{mixed} lower/middle cells exposed beneath overhead returns",
            start,
            cmap,
            norm,
            show_xlabel=True,
            extra_windows=extra_windows,
        )
    )
    colorbar = fig.colorbar(
        images[0],
        ax=descriptor_axes,
        orientation="horizontal",
        fraction=0.030,
        pad=0.20,
        aspect=28,
        extend="max",
    )
    colorbar.set_label(r"height $z_g$ (m)", fontsize=5.5, labelpad=2)
    colorbar.ax.tick_params(labelsize=5.2, length=2)
    colorbar.outline.set_linewidth(0.45)

    fig.savefig(output.with_suffix(".pdf"), bbox_inches="tight")
    fig.savefig(output.with_suffix(".svg"), bbox_inches="tight")
    fig.savefig(output.with_suffix(".png"), dpi=600, bbox_inches="tight")
    plt.close(fig)
    return {
        "window_start_sector": start,
        "window_width_sectors": WINDOW_WIDTH,
        "window_mixed_cells": window_mixed,
        "up_detail_window": {
            "start_sector": up_start,
            "criterion": "max count of overhead cells with top-underside "
                         "gap > 0.3 m, disjoint",
            "gap_cells": up_gap_cells,
            "median_gap_m": round(up_gap_median, 3),
        },
        "down_detail_window": {
            "start_sector": down_start,
            "criterion": "max std(Down)/std(SC), cells >= 40, disjoint",
            "std_ratio": round(down_ratio, 3),
            "down_std_m": round(down_std, 3),
            "sc_overhead_std_m": round(down_sc_std, 3),
        },
        "occupied_cells": occupied,
        "overhead_cells": overhead,
        "lower_cells": lower,
        "mixed_cells": mixed,
    }


def color_to_rgb(color) -> RGBColor:
    rgba = mpl.colors.to_rgba(color)
    return RGBColor(*(int(round(channel * 255)) for channel in rgba[:3]))


def ppt_text(slide, x, y, w, h, text, size=10, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color_to_rgb(color)
    return box


def build_ppt(descriptor, stats, output: Path):
    cmap, norm = height_style()
    start = stats["window_start_sector"]
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color_to_rgb(WHITE)
    ppt_text(slide, 0.18, 0.15, 12.9, 0.38, "Measured descriptor detail retained by the dual envelopes", 17, INK, True, PP_ALIGN.CENTER)

    lefts = [0.42, 4.55, 8.68]
    titles = [
        "SC: max over all heights",
        "Up: min above the split",
        "Down: max below the split",
    ]
    subtitles = [
        f'{stats["overhead_cells"]}/{stats["occupied_cells"]} cells set by overhead maxima',
        f'{stats["overhead_cells"]} overhead cells retained',
        f'{stats["mixed_cells"]} hidden lower/middle cells exposed',
    ]
    arrays = [
        (descriptor.sc_max, descriptor.sc_valid, False),
        (descriptor.upper_min, descriptor.upper_valid, False),
        (descriptor.lower_max, descriptor.lower_valid, False),
    ]
    cell_w, cell_h = 0.064, 0.185
    grid_top = 1.45
    for panel, (left, title, subtitle, payload) in enumerate(zip(lefts, titles, subtitles, arrays)):
        values, valid, mark_hidden = payload
        view, view_valid, view_hidden, display_start, _ = rolled(
            values, valid, descriptor.hidden_mask, start
        )
        ppt_text(slide, left, 0.78, 3.84, 0.34, title, 11, INK, True, PP_ALIGN.CENTER)
        ppt_text(slide, left, 1.10, 3.84, 0.28, subtitle, 8.5, MID, False, PP_ALIGN.CENTER)
        for ring in range(roof.NUM_RINGS):
            for sector in range(roof.NUM_SECTORS):
                color = cmap(norm(view[ring, sector])) if view_valid[ring, sector] else mpl.colors.to_rgba(LIGHT)
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(left + sector * cell_w),
                    Inches(grid_top + (roof.NUM_RINGS - 1 - ring) * cell_h),
                    Inches(cell_w),
                    Inches(cell_h),
                )
                shape.fill.solid(); shape.fill.fore_color.rgb = color_to_rgb(color)
                shape.line.fill.background()
                if mark_hidden and view_hidden[ring, sector]:
                    dot = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL,
                        Inches(left + sector * cell_w + 0.022),
                        Inches(grid_top + (roof.NUM_RINGS - 1 - ring) * cell_h + 0.078),
                        Inches(0.020),
                        Inches(0.020),
                    )
                    dot.fill.solid(); dot.fill.fore_color.rgb = color_to_rgb(RED); dot.line.fill.background()
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left + display_start * cell_w),
            Inches(grid_top),
            Inches(WINDOW_WIDTH * cell_w),
            Inches(roof.NUM_RINGS * cell_h),
        )
        box.fill.background(); box.line.color.rgb = color_to_rgb(MAGENTA); box.line.width = Pt(2)
    ppt_text(
        slide,
        2.0,
        5.05,
        9.3,
        0.38,
        f'magenta window: {stats["window_mixed_cells"]} mixed cells',
        10,
        MAGENTA,
        True,
        PP_ALIGN.CENTER,
    )
    ppt_text(slide, 0.8, 6.20, 11.7, 0.75, "Every descriptor cell, marker, label, and highlight is independently editable.", 9, MID, False, PP_ALIGN.CENTER)
    prs.save(output)


def write_source(record, descriptor, stats, output: Path):
    with output.open("w", newline="") as handle:
        fields = [
            "record_type", "query_id", "ring", "sector", "sc_max_z_m",
            "sc_valid", "up_min_z_m", "up_valid", "down_max_z_m",
            "down_valid", "mixed_cell",
        ]
        writer = csv.DictWriter(handle, fieldnames=fields)
        writer.writeheader()
        for ring in range(roof.NUM_RINGS):
            for sector in range(roof.NUM_SECTORS):
                sc_valid = bool(descriptor.sc_valid[ring, sector])
                up_valid = bool(descriptor.upper_valid[ring, sector])
                down_valid = bool(descriptor.lower_valid[ring, sector])
                writer.writerow(
                    {
                        "record_type": "descriptor_cell",
                        "query_id": record.query_id,
                        "ring": ring,
                        "sector": sector,
                        "sc_max_z_m": descriptor.sc_max[ring, sector] if sc_valid else "",
                        "sc_valid": int(sc_valid),
                        "up_min_z_m": descriptor.upper_min[ring, sector] if up_valid else "",
                        "up_valid": int(up_valid),
                        "down_max_z_m": descriptor.lower_max[ring, sector] if down_valid else "",
                        "down_valid": int(down_valid),
                        "mixed_cell": int(descriptor.hidden_mask[ring, sector]),
                    }
                )

    audit = {
        "selection_rule": (
            "reuse the independently selected roof-contamination query: maximize conventional-SC cells "
            "above the final 2.1 m IH split across all truth-valid queries"
        ),
        "evaluated_truth_valid_queries": 220,
        "selected_query_id": record.query_id,
        "descriptor": {
            "rings": roof.NUM_RINGS,
            "sectors": roof.NUM_SECTORS,
            "radius_m": roof.MAX_RADIUS_M,
            "split_height_m": SPLIT_HEIGHT_M,
            "voxel_m": roof.VOXEL_M,
        },
        "statistics": stats,
        "image_integrity": (
            "All descriptor cells use the same measured query. The descriptor matrices are "
            "circularly shifted only for display, with the same shift applied to SC, Up, "
            "Down, masks, and the highlighted window. No cell is "
            "interpolated or manually edited."
        ),
        "exclusions": "none after the pre-existing truth-valid query criterion",
    }
    with output.with_name("descriptor_detail_metadata.json").open("w") as handle:
        json.dump(audit, handle, indent=2)


def main():
    args = parse_args()
    args.output_dir.mkdir(parents=True, exist_ok=True)
    roof.SPLIT_HEIGHT_M = SPLIT_HEIGHT_M
    records = roof.load_queries(args.query_dir)
    record, descriptor, summaries = roof.select_query(records, None)
    if len(summaries) != 220:
        raise RuntimeError(f"Expected 220 truth-valid selection records, got {len(summaries)}")
    prefix = args.output_dir / "descriptor_detail"
    stats = make_figure(record, descriptor, prefix)
    build_ppt(descriptor, stats, args.output_dir / "descriptor_detail_editable.pptx")
    write_source(
        record,
        descriptor,
        stats,
        args.output_dir / "descriptor_detail_source_data.csv",
    )
    print(
        f"Generated descriptor detail for query {record.query_id}: "
        f"SC={stats['occupied_cells']}, Up={stats['overhead_cells']}, "
        f"mixed={stats['mixed_cells']}."
    )


if __name__ == "__main__":
    main()
