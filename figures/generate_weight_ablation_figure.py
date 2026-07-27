#!/usr/bin/env python3
"""Generate the real-data UpDown-SC channel-weight sensitivity figure.

Figure contract
---------------
Core conclusion:
    The (lower, upper) weight pair (0.3, 0.7), selected on the IH+G pilot and
    frozen elsewhere, also maximizes unweighted macro Recall@1; the full sweep
    exposes real per-condition trade-offs rather than a uniform gain.
Evidence chain:
    Panel (a) shows macro Recall@1/5 for all 11 tested pairs. Panel (b) shows
    every condition/pair Recall@1 change relative to equal weighting.
Archetype:
    Quantitative grid with a summary hero panel and a complete sensitivity map.
Outputs:
    PDF/SVG/TIFF/PNG, native-shape editable PowerPoint, source CSV, and metadata.
"""

from __future__ import annotations

import argparse
import csv
import json
from collections import defaultdict
from pathlib import Path

import matplotlib as mpl
import matplotlib.pyplot as plt
import numpy as np
from matplotlib.colors import LinearSegmentedColormap, TwoSlopeNorm
from matplotlib.patches import Rectangle
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parent
DEFAULT_INPUT = (
    Path.home()
    / "icra2027_runtime/experiments/updown_weight_ablation_real_20260721"
    / "weight_sweep_by_dataset.csv"
)

DATASETS = [
    ("ih_native", "IH native"),
    ("ih_gravity", "IH +G"),
    ("ch_native", "CH native"),
    ("ch_gravity", "CH +G"),
    ("h1_h2", "H1→H2 +G"),
    ("h1_v1", "H1→V1 +G"),
    ("nc_gravity", "NC +G"),
]
SELECTED_WEIGHT = 0.3
EQUAL_WEIGHT = 0.5

BLUE = "#2468A9"
ORANGE = "#D88318"
TEAL = "#218C82"
INK = "#1B2632"
MID = "#687483"
GRID = "#D7DEE5"
WHITE = "#FFFFFF"
HEAT_NEG = "#4E78A8"
HEAT_ZERO = "#F5F5F2"
HEAT_POS = "#D36B4A"

mpl.rcParams.update(
    {
        "font.family": "sans-serif",
        "font.sans-serif": ["Arial", "Helvetica", "DejaVu Sans", "sans-serif"],
        "font.size": 6.4,
        "axes.labelsize": 6.6,
        "axes.titlesize": 7.2,
        "axes.linewidth": 0.7,
        "xtick.labelsize": 5.8,
        "ytick.labelsize": 5.8,
        "legend.fontsize": 6.1,
        "legend.frameon": False,
        "svg.fonttype": "none",
        "pdf.fonttype": 42,
    }
)


def load_rows(path: Path):
    with path.open(newline="") as handle:
        raw = list(csv.DictReader(handle))
    required = {
        "lower_weight",
        "upper_weight",
        "dataset",
        "label",
        "queries",
        "recall_at_1",
        "recall_at_5",
    }
    if not raw or not required.issubset(raw[0]):
        raise RuntimeError(f"Unexpected weight-sweep schema: {path}")

    parsed = []
    for row in raw:
        parsed.append(
            {
                "lower_weight": round(float(row["lower_weight"]), 1),
                "upper_weight": round(float(row["upper_weight"]), 1),
                "dataset": row["dataset"],
                "label": row["label"],
                "queries": int(row["queries"]),
                "recall_at_1": 100.0 * float(row["recall_at_1"]),
                "recall_at_5": 100.0 * float(row["recall_at_5"]),
            }
        )

    weights = sorted({row["lower_weight"] for row in parsed})
    expected_weights = [round(value / 10.0, 1) for value in range(11)]
    if weights != expected_weights:
        raise RuntimeError(f"Expected weights 0.0,...,1.0, got {weights}")
    expected_datasets = {key for key, _ in DATASETS}
    if {row["dataset"] for row in parsed} != expected_datasets:
        raise RuntimeError("Dataset set differs from the seven Table-II conditions")
    if len(parsed) != len(weights) * len(DATASETS):
        raise RuntimeError(f"Expected 77 sweep records, got {len(parsed)}")
    return parsed, weights


def assemble(rows, weights):
    by_dataset = defaultdict(dict)
    for row in rows:
        by_dataset[row["dataset"]][row["lower_weight"]] = row

    r1 = np.array(
        [[by_dataset[key][weight]["recall_at_1"] for weight in weights]
         for key, _ in DATASETS]
    )
    r5 = np.array(
        [[by_dataset[key][weight]["recall_at_5"] for weight in weights]
         for key, _ in DATASETS]
    )
    macro_r1 = np.mean(r1, axis=0)
    macro_r5 = np.mean(r5, axis=0)
    equal_index = weights.index(EQUAL_WEIGHT)
    delta_r1 = r1 - r1[:, [equal_index]]
    selected_index = weights.index(SELECTED_WEIGHT)
    pilot_index = [key for key, _ in DATASETS].index("ih_gravity")
    best_index = int(np.argmax(r1[pilot_index]))
    if best_index != selected_index:
        raise RuntimeError(
            f"Selected weight {SELECTED_WEIGHT:.1f} is not the IH+G-pilot R@1 maximum"
        )
    return by_dataset, r1, r5, macro_r1, macro_r5, delta_r1, selected_index


def draw_curve_panel(ax, weights, macro_r1, macro_r5, selected_index):
    ax.plot(
        weights,
        macro_r1,
        color=BLUE,
        marker="o",
        markersize=3.1,
        linewidth=1.55,
        label="Macro R@1",
        zorder=3,
    )
    ax.plot(
        weights,
        macro_r5,
        color=ORANGE,
        marker="s",
        markersize=2.8,
        linewidth=1.45,
        label="Macro R@5",
        zorder=3,
    )
    ax.axvline(SELECTED_WEIGHT, color=TEAL, linewidth=1.0, linestyle=(0, (3, 2)))
    ax.scatter(
        [SELECTED_WEIGHT],
        [macro_r1[selected_index]],
        s=34,
        marker="*",
        color=TEAL,
        edgecolor=WHITE,
        linewidth=0.5,
        zorder=5,
    )
    ax.annotate(
        "pilot-selected 0.3/0.7",
        (SELECTED_WEIGHT, macro_r1[selected_index]),
        xytext=(5, 8),
        textcoords="offset points",
        color=TEAL,
        fontsize=5.8,
        fontweight="bold",
    )
    ax.set_xlim(-0.025, 1.025)
    ax.set_ylim(52, 86)
    ax.set_xticks(np.arange(0.0, 1.01, 0.2))
    ax.set_ylabel("Macro recall (%)")
    ax.set_xlabel(r"Down (lower-layer) weight $w_\ell$  ($w_h=1-w_\ell$)")
    ax.grid(axis="y", color=GRID, linewidth=0.55)
    ax.spines[["top", "right"]].set_visible(False)
    ax.legend(loc="lower center", bbox_to_anchor=(0.52, 1.01), ncol=2)


def draw_heatmap_panel(axh, weights, delta_r1, selected_index):
    cmap = LinearSegmentedColormap.from_list(
        "weight_delta", [HEAT_NEG, HEAT_ZERO, HEAT_POS]
    )
    norm = TwoSlopeNorm(vmin=-12.0, vcenter=0.0, vmax=12.0)
    image = axh.imshow(delta_r1, aspect="auto", cmap=cmap, norm=norm)
    axh.set_yticks(np.arange(len(DATASETS)), [label for _, label in DATASETS])
    axh.set_xticks(np.arange(len(weights)))
    axh.set_xticklabels([f"{weight:.1f}" for weight in weights], rotation=0)
    axh.set_xlabel(r"Down (lower-layer) weight $w_\ell$")
    axh.tick_params(length=0)
    for row in range(len(DATASETS)):
        best = np.flatnonzero(np.isclose(delta_r1[row], np.max(delta_r1[row])))
        for column in best:
            axh.scatter(column, row, s=8, marker="o", facecolor="none",
                        edgecolor=INK, linewidth=0.55)
    axh.add_patch(
        Rectangle(
            (selected_index - 0.5, -0.5),
            1.0,
            len(DATASETS),
            fill=False,
            edgecolor=TEAL,
            linewidth=1.4,
        )
    )
    for x in np.arange(-0.5, len(weights), 1.0):
        axh.axvline(x, color=WHITE, linewidth=0.35)
    for y in np.arange(-0.5, len(DATASETS), 1.0):
        axh.axhline(y, color=WHITE, linewidth=0.35)
    for spine in axh.spines.values():
        spine.set_visible(False)
    return image


def add_heatmap_colorbar(fig, axh, image):
    cbar = fig.colorbar(image, ax=axh, orientation="horizontal", fraction=0.075, pad=0.18)
    cbar.set_label("R@1 change from equal weighting (percentage points)", labelpad=2)
    cbar.set_ticks([-12, -6, 0, 6, 12])
    cbar.outline.set_linewidth(0.45)
    cbar.ax.tick_params(labelsize=5.5, length=2)


def save_figure(fig, output: Path, include_tiff: bool = False):
    fig.savefig(output.with_suffix(".pdf"), bbox_inches="tight")
    fig.savefig(output.with_suffix(".svg"), bbox_inches="tight")
    fig.savefig(output.with_suffix(".png"), dpi=600, bbox_inches="tight")
    if include_tiff:
        fig.savefig(output.with_suffix(".tiff"), dpi=600, bbox_inches="tight")
    plt.close(fig)


def make_figure(weights, macro_r1, macro_r5, delta_r1, selected_index, output: Path):
    # Legacy combined export, now intentionally free of embedded panel letters.
    fig = plt.figure(figsize=(3.45, 4.05), constrained_layout=False)
    grid = fig.add_gridspec(2, 1, height_ratios=[1.0, 1.42], hspace=0.58)
    ax = fig.add_subplot(grid[0])
    draw_curve_panel(ax, weights, macro_r1, macro_r5, selected_index)
    axh = fig.add_subplot(grid[1])
    image = draw_heatmap_panel(axh, weights, delta_r1, selected_index)
    add_heatmap_colorbar(fig, axh, image)
    fig.subplots_adjust(left=0.25, right=0.985, top=0.955, bottom=0.10)
    save_figure(fig, output, include_tiff=True)

    # Manuscript panels are exported separately; LaTeX owns the (a)/(b) labels.
    curve_fig, curve_ax = plt.subplots(figsize=(3.45, 1.55), facecolor=WHITE)
    draw_curve_panel(curve_ax, weights, macro_r1, macro_r5, selected_index)
    curve_fig.subplots_adjust(left=0.16, right=0.985, top=0.80, bottom=0.27)
    save_figure(curve_fig, output.with_name(output.name + "_curve"))

    heat_fig, heat_ax = plt.subplots(figsize=(3.45, 2.15), facecolor=WHITE)
    heat_image = draw_heatmap_panel(heat_ax, weights, delta_r1, selected_index)
    add_heatmap_colorbar(heat_fig, heat_ax, heat_image)
    heat_fig.subplots_adjust(left=0.25, right=0.985, top=0.97, bottom=0.23)
    save_figure(heat_fig, output.with_name(output.name + "_heatmap"))


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_text(slide, x, y, w, h, text, size=12, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def build_editable_ppt(weights, macro_r1, macro_r5, delta_r1, selected_index, output: Path):
    prs = Presentation()
    prs.slide_width = Inches(7.5)
    prs.slide_height = Inches(8.8)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(WHITE)

    left, top, width, height = 1.10, 0.72, 5.95, 2.45
    ymin, ymax = 52.0, 86.0
    for value in [60, 70, 80]:
        py = top + height * (ymax - value) / (ymax - ymin)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(py), Inches(width), Inches(0.008))
        line.fill.solid(); line.fill.fore_color.rgb = rgb(GRID); line.line.fill.background()
        add_text(slide, 0.42, py - 0.12, 0.55, 0.25, str(value), 9, align=PP_ALIGN.RIGHT)
    axis = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.012), Inches(height))
    axis.fill.solid(); axis.fill.fore_color.rgb = rgb(INK); axis.line.fill.background()
    axis = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top + height), Inches(width), Inches(0.012))
    axis.fill.solid(); axis.fill.fore_color.rgb = rgb(INK); axis.line.fill.background()

    def point(weight, value):
        return (
            left + width * weight,
            top + height * (ymax - value) / (ymax - ymin),
        )

    for values, color, marker in [(macro_r1, BLUE, MSO_SHAPE.OVAL), (macro_r5, ORANGE, MSO_SHAPE.RECTANGLE)]:
        previous = None
        for weight, value in zip(weights, values):
            px, py = point(weight, value)
            if previous is not None:
                connector = slide.shapes.add_connector(1, Inches(previous[0]), Inches(previous[1]), Inches(px), Inches(py))
                connector.line.color.rgb = rgb(color); connector.line.width = Pt(2)
            shape = slide.shapes.add_shape(marker, Inches(px - 0.055), Inches(py - 0.055), Inches(0.11), Inches(0.11))
            shape.fill.solid(); shape.fill.fore_color.rgb = rgb(color); shape.line.fill.background()
            previous = (px, py)
    sx, sy = point(SELECTED_WEIGHT, macro_r1[selected_index])
    guide = slide.shapes.add_connector(1, Inches(sx), Inches(top), Inches(sx), Inches(top + height))
    guide.line.color.rgb = rgb(TEAL); guide.line.width = Pt(1.3); guide.line.dash_style = 2
    add_text(slide, sx + 0.08, sy - 0.34, 1.65, 0.3, "pilot-selected 0.3/0.7", 9, TEAL, True)
    for weight in np.arange(0, 1.01, 0.2):
        px, _ = point(float(weight), ymin)
        add_text(slide, px - 0.20, top + height + 0.08, 0.40, 0.25, f"{weight:.1f}", 8, align=PP_ALIGN.CENTER)
    add_text(slide, 2.23, 3.43, 3.7, 0.35, "Down (lower-layer) weight w_l  (w_h = 1 - w_l)", 10, align=PP_ALIGN.CENTER)
    add_text(slide, 2.03, 0.20, 1.55, 0.3, "●  Macro R@1", 10, BLUE, True)
    add_text(slide, 4.02, 0.20, 1.55, 0.3, "■  Macro R@5", 10, ORANGE, True)

    heat_left, heat_top = 1.42, 4.38
    cell_w, cell_h = 0.49, 0.45
    norm = TwoSlopeNorm(vmin=-12.0, vcenter=0.0, vmax=12.0)
    cmap = LinearSegmentedColormap.from_list("weight_delta", [HEAT_NEG, HEAT_ZERO, HEAT_POS])
    for row, (_, label) in enumerate(DATASETS):
        add_text(slide, 0.15, heat_top + row * cell_h + 0.10, 1.15, 0.22, label, 8.5, align=PP_ALIGN.RIGHT)
        row_best = np.max(delta_r1[row])
        for col, value in enumerate(delta_r1[row]):
            color = mpl.colors.to_hex(cmap(norm(value)))
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(heat_left + col * cell_w),
                Inches(heat_top + row * cell_h),
                Inches(cell_w),
                Inches(cell_h),
            )
            shape.fill.solid(); shape.fill.fore_color.rgb = rgb(color)
            shape.line.color.rgb = rgb(WHITE); shape.line.width = Pt(0.4)
            if np.isclose(value, row_best):
                dot = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(heat_left + col * cell_w + 0.205),
                    Inches(heat_top + row * cell_h + 0.185),
                    Inches(0.08),
                    Inches(0.08),
                )
                dot.fill.background(); dot.line.color.rgb = rgb(INK); dot.line.width = Pt(0.8)
    border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(heat_left + selected_index * cell_w),
        Inches(heat_top),
        Inches(cell_w),
        Inches(cell_h * len(DATASETS)),
    )
    border.fill.background(); border.line.color.rgb = rgb(TEAL); border.line.width = Pt(2)
    for col, weight in enumerate(weights):
        add_text(slide, heat_left + col * cell_w, heat_top + cell_h * len(DATASETS) + 0.08, cell_w, 0.22, f"{weight:.1f}", 7.5, align=PP_ALIGN.CENTER)
    add_text(slide, 2.2, 7.82, 3.8, 0.3, "Down (lower-layer) weight w_l", 10, align=PP_ALIGN.CENTER)
    add_text(slide, 1.55, 8.23, 5.4, 0.3, "Color: R@1 change from equal weighting (percentage points)", 9, MID, align=PP_ALIGN.CENTER)
    prs.save(output)


def export_source(rows, weights, r1, r5, delta_r1, macro_r1, macro_r5, output: Path):
    with output.open("w", newline="") as handle:
        fields = [
            "record_type", "dataset", "label", "queries", "lower_weight",
            "upper_weight", "recall_at_1_percent", "recall_at_5_percent",
            "delta_r1_from_equal_pp",
        ]
        writer = csv.DictWriter(handle, fieldnames=fields)
        writer.writeheader()
        for row_index, (dataset, label) in enumerate(DATASETS):
            query_count = next(row["queries"] for row in rows if row["dataset"] == dataset)
            for column, weight in enumerate(weights):
                writer.writerow(
                    {
                        "record_type": "condition",
                        "dataset": dataset,
                        "label": label,
                        "queries": query_count,
                        "lower_weight": f"{weight:.1f}",
                        "upper_weight": f"{1.0 - weight:.1f}",
                        "recall_at_1_percent": f"{r1[row_index, column]:.9f}",
                        "recall_at_5_percent": f"{r5[row_index, column]:.9f}",
                        "delta_r1_from_equal_pp": f"{delta_r1[row_index, column]:.9f}",
                    }
                )
        for column, weight in enumerate(weights):
            writer.writerow(
                {
                    "record_type": "macro",
                    "dataset": "macro_unweighted",
                    "label": "Seven-condition macro average",
                    "queries": "",
                    "lower_weight": f"{weight:.1f}",
                    "upper_weight": f"{1.0 - weight:.1f}",
                    "recall_at_1_percent": f"{macro_r1[column]:.9f}",
                    "recall_at_5_percent": f"{macro_r5[column]:.9f}",
                    "delta_r1_from_equal_pp": f"{macro_r1[column] - macro_r1[weights.index(EQUAL_WEIGHT)]:.9f}",
                }
            )


def write_metadata(output: Path):
    metadata = {
        "figure_contract": {
            "core_conclusion": (
                "The IH+G-pilot-selected 0.3/0.7 pair also maximizes macro "
                "Recall@1, while every condition-level response remains visible."
            ),
            "archetype": "quantitative grid",
            "backend": "Python/Matplotlib with native-shape PowerPoint export",
            "target_width_mm": 87.6,
        },
        "input": "archived real-data weight-sweep CSV",
        "records": 77,
        "conditions": 7,
        "weight_pairs": 11,
        "selection": "maximize IH+G pilot Recall@1, then freeze for every other condition",
        "baseline_for_heatmap": "equal weighting (0.5, 0.5)",
        "exclusions": "none",
        "uncertainty": "deterministic retrieval on fixed queries; no seed variance claimed",
    }
    with output.open("w") as handle:
        json.dump(metadata, handle, indent=2)


def parse_args():
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", type=Path, default=DEFAULT_INPUT)
    parser.add_argument("--output-dir", type=Path, default=OUT)
    return parser.parse_args()


def main():
    args = parse_args()
    args.output_dir.mkdir(parents=True, exist_ok=True)
    rows, weights = load_rows(args.input)
    _, r1, r5, macro_r1, macro_r5, delta_r1, selected_index = assemble(rows, weights)
    prefix = args.output_dir / "weight_ablation"
    make_figure(weights, macro_r1, macro_r5, delta_r1, selected_index, prefix)
    build_editable_ppt(
        weights,
        macro_r1,
        macro_r5,
        delta_r1,
        selected_index,
        args.output_dir / "weight_ablation_editable.pptx",
    )
    export_source(
        rows,
        weights,
        r1,
        r5,
        delta_r1,
        macro_r1,
        macro_r5,
        args.output_dir / "weight_ablation_source_data.csv",
    )
    write_metadata(args.output_dir / "weight_ablation_metadata.json")
    print("Generated weight-ablation figure from 77 real-data sweep records.")


if __name__ == "__main__":
    main()
