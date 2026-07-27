#!/usr/bin/env python3
"""Generate editable PowerPoint sources for method figures.

Text, grids, polar cells, arrows, and colorbar steps are independent PowerPoint
objects. The measured cloud is the sole raster layer because tens of thousands
of point markers are not practical as PowerPoint shapes.
"""

from __future__ import annotations

import argparse
import math
from pathlib import Path

import numpy as np
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import generate_editable_ppt as common
import generate_roof_contamination_figure as roof


OUT_DIR = Path(__file__).resolve().parent
PRINCIPLE_OUT = OUT_DIR / "updown_sc_principle_editable.pptx"
ROOF_OUT = OUT_DIR / "roof_contamination_editable.pptx"


def blank_slide(width: float, height: float):
    prs = Presentation()
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = common.rgb(common.WHITE)
    return prs, slide


def add_panel(slide, x: float, y: float, w: float, h: float, label: str, title: str, fill: str):
    common.add_rect(slide, x, y, w, h, fill=fill, line=common.BORDER, line_width=0.8)
    common.add_text(slide, x + 0.08, y + 0.06, 0.24, 0.24, label, 10, common.INK, True, PP_ALIGN.LEFT)
    common.add_text(slide, x + 0.34, y + 0.06, w - 0.42, 0.24, title, 8.4, common.INK, True, PP_ALIGN.LEFT)


def draw_room(slide, x: float, y: float, w: float, h: float, dual: bool):
    def px(v: float) -> float:
        return x + v * w

    def py(v: float) -> float:
        return y + (1.0 - v) * h

    common.add_line(slide, px(0.12), py(0.78), px(0.90), py(0.78), common.INK, 1.6)
    common.add_text(slide, px(0.69), py(0.89), 0.60, 0.18, "ceiling", 6.2, common.INK)
    common.add_rect(slide, px(0.19), py(0.48), 0.22 * w, 0.28 * h, common.BLUE_LIGHT, common.BLUE_2, 0.7)
    common.add_text(slide, px(0.20), py(0.20) + 0.02, 0.62, 0.18, "shelf", 6.0, common.BLUE)
    common.add_rect(slide, px(0.56), py(0.58), 0.08 * w, 0.38 * h, "#E6E8EB", common.MID, 0.7)
    common.add_oval(slide, px(0.555), py(0.675), 0.09 * w, 0.09 * w, "#E6E8EB", common.MID, 0.7)
    common.add_text(slide, px(0.49), py(0.17) + 0.02, 0.62, 0.18, "person", 6.0, common.MID)
    common.add_line(slide, px(0.10), py(0.20), px(0.92), py(0.20), common.MID, 0.7)
    for xx, yy in zip((0.22, 0.29, 0.37, 0.58, 0.62, 0.71, 0.79), (0.48, 0.47, 0.48, 0.58, 0.56, 0.78, 0.78)):
        common.add_oval(slide, px(xx) - 0.028, py(yy) - 0.028, 0.056, 0.056, common.INK, common.INK, 0)

    if not dual:
        common.add_line(slide, px(0.13), py(0.78), px(0.88), py(0.78), "#B64342", 2.4)
        common.add_text(slide, px(0.30), py(0.70), 1.50, 0.20, "selected max z", 6.4, "#B64342", True)
        common.add_line(slide, px(0.21), py(0.42), px(0.40), py(0.24), "#B64342", 0.9)
        common.add_line(slide, px(0.40), py(0.42), px(0.21), py(0.24), "#B64342", 0.9)
        common.add_text(slide, px(0.12), py(0.12), 1.25, 0.34, "lower structure\nnot represented", 5.8, "#B64342")
        return

    threshold = common.add_line(slide, px(0.12), py(0.59), px(0.90), py(0.59), common.ORANGE, 0.9, True)
    threshold.line.dash_style = common.MSO_LINE_DASH_STYLE.DASH
    common.add_text(slide, px(0.89), py(0.62), 0.18, 0.20, "τ", 7.0, common.ORANGE)
    common.add_line(slide, px(0.18), py(0.48), px(0.42), py(0.48), common.BLUE, 2.4)
    common.add_line(slide, px(0.66), py(0.78), px(0.87), py(0.78), common.TEAL, 2.4)
    common.add_text(slide, px(0.09), py(0.59), 1.50, 0.20, "Down: max below τ", 5.8, common.BLUE, True)
    common.add_text(slide, px(0.55), py(0.74), 1.30, 0.20, "Up: min above τ", 5.8, common.TEAL, True)


def binary_grid(slide, x: float, y: float, values: np.ndarray, cell: float, color: str, label: str):
    rows, cols = values.shape
    common.add_text(slide, x, y - 0.22, cols * cell, 0.18, label, 6.1, common.INK)
    for row in range(rows):
        for col in range(cols):
            value = int(values[row, col])
            common.add_rect(
                slide,
                x + col * cell,
                y + row * cell,
                cell * 0.92,
                cell * 0.92,
                fill=color if value else common.WHITE,
                line=common.BORDER,
                line_width=0.35,
            )


def build_principle() -> None:
    prs, slide = blank_slide(7.0, 7.56)
    common.add_text(slide, 0.20, 0.06, 6.60, 0.34, "Dual-envelope encoding and matching", 15, common.INK, True)
    boxes = ((0.10, 0.55), (3.55, 0.55), (0.10, 3.90), (3.55, 3.90))
    for (x, y), label, title, fill in zip(
        boxes,
        "abcd",
        ("Conventional SC", "Strict dual envelope", "Joint-valid masks", "Conditional vertical seed"),
        ("#FCF8F8", "#F7FAFD", "#FAFBFC", "#FAFBFC"),
    ):
        add_panel(slide, x, y, 3.35, 3.15, label, title, fill)

    draw_room(slide, 0.30, 1.02, 2.90, 2.38, False)
    draw_room(slide, 3.75, 1.02, 2.90, 2.38, True)

    query = np.array([[1, 1, 0, 1, 0], [1, 0, 1, 1, 1], [0, 1, 1, 0, 1]], dtype=int)
    mapped = np.array([[1, 0, 1, 1, 0], [1, 1, 1, 0, 1], [0, 1, 0, 1, 1]], dtype=int)
    joint = query & mapped
    binary_grid(slide, 0.35, 4.65, query, 0.18, common.BLUE_2, "query M")
    common.add_text(slide, 1.28, 4.78, 0.28, 0.28, "∧", 13, common.MID)
    binary_grid(slide, 1.62, 4.65, mapped, 0.18, common.TEAL, "map M")
    common.add_right_arrow(slide, 2.55, 4.87, 0.28, 0.12, common.MID)
    binary_grid(slide, 2.05, 5.55, joint, 0.18, common.PURPLE, "joint m")
    common.add_text(slide, 0.35, 5.55, 1.48, 0.42, "cosine uses only\njoint-valid rings", 6.4, common.INK)
    common.add_text(slide, 0.45, 6.08, 1.28, 0.20, "joint rings ≥ 2", 6.2, common.PURPLE)
    common.add_tag(slide, 0.58, 6.56, 2.38, 0.30, "mask similarity × sector support", common.PURPLE_LIGHT, common.PURPLE, 6.6)

    x, y = 3.55, 3.90
    common.add_text(slide, x + 0.75, y + 0.52, 1.85, 0.22, "place + yaw fixed", 6.5, common.PURPLE, True)
    common.add_line(slide, x + 0.45, y + 1.22, x + 2.93, y + 1.22, common.ORANGE, 0.9, True)
    common.add_text(slide, x + 2.93, y + 1.11, 0.20, 0.20, "τ", 7, common.ORANGE)
    q_lower = (0.34, 0.42, 0.50, 0.55)
    m_lower = (0.38, 0.47, 0.53, 0.57)
    q_upper = (0.68, 0.73, 0.78)
    m_upper = (0.72, 0.76, 0.82)
    for q, m in zip(q_lower, m_lower):
        yq, ym = y + 2.60 - q * 2.0, y + 2.60 - m * 2.0
        common.add_line(slide, x + 0.62, yq, x + 1.28, ym, common.BLUE_2, 0.8)
        common.add_oval(slide, x + 0.59, yq - 0.03, 0.06, 0.06, common.BLUE_2, common.BLUE_2, 0)
        common.add_oval(slide, x + 1.25, ym - 0.03, 0.06, 0.06, common.BLUE_2, common.BLUE_2, 0)
    for q, m in zip(q_upper, m_upper):
        yq, ym = y + 2.60 - q * 2.0, y + 2.60 - m * 2.0
        common.add_line(slide, x + 0.62, yq, x + 1.28, ym, common.TEAL, 0.8)
        common.add_oval(slide, x + 0.59, yq - 0.03, 0.06, 0.06, common.TEAL, common.TEAL, 0)
        common.add_oval(slide, x + 1.25, ym - 0.03, 0.06, 0.06, common.TEAL, common.TEAL, 0)
    common.add_text(slide, x + 0.36, y + 2.43, 0.54, 0.18, "query", 5.8, common.MID)
    common.add_text(slide, x + 1.03, y + 2.43, 0.54, 0.18, "map", 5.8, common.MID)
    common.add_right_arrow(slide, x + 1.50, y + 1.70, 0.34, 0.14, common.ORANGE)
    common.add_text(slide, x + 1.38, y + 1.42, 0.62, 0.20, "r = zm − zq", 5.8, common.ORANGE)
    common.add_text(slide, x + 1.98, y + 0.98, 1.15, 0.38, "rank by distance\nfrom split τ", 5.8, common.INK)
    common.add_text(slide, x + 1.98, y + 1.70, 1.15, 0.24, "keep stable 50%", 6.2, common.PURPLE, True)
    common.add_right_arrow(slide, x + 2.43, y + 2.03, 0.15, 0.28, common.ORANGE, 90)
    common.add_tag(slide, x + 1.93, y + 2.50, 1.28, 0.29, "weighted median → Δz", "#FCE8D1", common.ORANGE, 5.8)
    common.add_text(slide, 2.15, 7.25, 4.65, 0.18, "Conceptual schematic; τ and 50% are configurable in the implementation", 5.4, common.MID, False, PP_ALIGN.RIGHT)
    prs.save(PRINCIPLE_OUT)


def interpolate_color(value: float) -> str:
    stops = np.array(
        [
            [0x26, 0x3B, 0x73],
            [0x37, 0x75, 0xBA],
            [0xD9, 0xEF, 0xEA],
            [0xF4, 0xC6, 0x6A],
            [0xB6, 0x43, 0x42],
        ],
        dtype=float,
    )
    t = float(np.clip((value + 0.5) / 4.5, 0.0, 1.0))
    scaled = t * (len(stops) - 1)
    idx = min(int(math.floor(scaled)), len(stops) - 2)
    frac = scaled - idx
    color = np.rint(stops[idx] * (1.0 - frac) + stops[idx + 1] * frac).astype(int)
    return "#" + "".join(f"{channel:02X}" for channel in color)


def add_wedge(slide, cx: float, cy: float, radius: float, ring: int, sector: int, color: str):
    inner = radius * ring / roof.NUM_RINGS
    outer = radius * (ring + 1) / roof.NUM_RINGS
    a0 = sector * 2.0 * math.pi / roof.NUM_SECTORS
    a1 = (sector + 1) * 2.0 * math.pi / roof.NUM_SECTORS
    angles_outer = np.linspace(a0, a1, 3)
    angles_inner = np.linspace(a1, a0, 3)
    vertices = []
    for angle in angles_outer:
        vertices.append(((cx + outer * math.cos(angle)) * 1000, (cy - outer * math.sin(angle)) * 1000))
    for angle in angles_inner:
        vertices.append(((cx + inner * math.cos(angle)) * 1000, (cy - inner * math.sin(angle)) * 1000))
    builder = slide.shapes.build_freeform(vertices[0][0], vertices[0][1], scale=Inches(1) / 1000)
    builder.add_line_segments(vertices[1:], close=True)
    shape = builder.convert_to_shape()
    shape.fill.solid()
    shape.fill.fore_color.rgb = common.rgb(color)
    shape.line.color.rgb = common.rgb(color)
    shape.line.width = Pt(0.1)
    return shape


def add_polar_descriptor(
    slide,
    cx: float,
    cy: float,
    radius: float,
    values: np.ndarray,
    valid: np.ndarray,
    title: str,
    hidden: np.ndarray | None = None,
):
    common.add_oval(slide, cx - radius, cy - radius, 2 * radius, 2 * radius, roof.LIGHT, roof.BORDER, 0.7)
    for ring, sector in zip(*np.nonzero(valid)):
        add_wedge(slide, cx, cy, radius, int(ring), int(sector), interpolate_color(float(values[ring, sector])))
    for frac in (1.0 / 3.0, 2.0 / 3.0):
        r = radius * frac
        common.add_oval(slide, cx - r, cy - r, 2 * r, 2 * r, None, common.WHITE, 0.45)
    for angle in (0.0, math.pi / 2.0, math.pi, 3.0 * math.pi / 2.0):
        common.add_line(slide, cx, cy, cx + radius * math.cos(angle), cy - radius * math.sin(angle), common.WHITE, 0.45)
    common.add_oval(slide, cx - radius, cy - radius, 2 * radius, 2 * radius, None, roof.BORDER, 0.75)
    if hidden is not None:
        for ring, sector in zip(*np.nonzero(hidden)):
            theta = (sector + 0.5) * 2.0 * math.pi / roof.NUM_SECTORS
            rr = (ring + 0.5) * radius / roof.NUM_RINGS
            common.add_oval(slide, cx + rr * math.cos(theta) - 0.025, cy - rr * math.sin(theta) - 0.025, 0.05, 0.05, roof.RED, roof.RED, 0)
    common.add_text(slide, cx - radius, cy - radius - 0.47, 2 * radius, 0.30, title, 11.5, roof.INK, True)


def add_colorbar(slide, x: float, y: float, w: float, h: float):
    steps = 24
    for idx in range(steps):
        value = -0.5 + (idx + 0.5) / steps * 4.5
        yy = y + h - (idx + 1) * h / steps
        common.add_rect(slide, x, yy, w, h / steps + 0.004, interpolate_color(value), interpolate_color(value), 0)
    triangle = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x), Inches(y - 0.17), Inches(w), Inches(0.18))
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = common.rgb(interpolate_color(4.0))
    triangle.line.color.rgb = common.rgb(common.INK)
    triangle.line.width = Pt(0.6)
    common.add_rect(slide, x, y, w, h, fill=common.WHITE, line=common.INK, line_width=0.7).fill.background()
    for value in np.arange(-0.5, 4.01, 0.5):
        yy = y + h - (value + 0.5) / 4.5 * h
        common.add_line(slide, x + w, yy, x + w + 0.07, yy, common.INK, 0.6)
        common.add_text(slide, x + w + 0.09, yy - 0.10, 0.48, 0.20, f"{value:.1f}", 6.8, common.INK, False, PP_ALIGN.LEFT)
    common.add_text(slide, x + 0.47, y + 0.58, 0.42, h - 0.70, "height z (m)", 8.5, common.INK, False, PP_ALIGN.CENTER).rotation = 270


def build_roof(query_dir: Path) -> None:
    records = roof.load_queries(query_dir)
    _, descriptor, _ = roof.select_query(records, None)
    prs, slide = blank_slide(13.333, 4.25)
    cloud_path = OUT_DIR / "roof_contamination_cloud.png"
    slide.shapes.add_picture(str(cloud_path), Inches(0.12), Inches(0.28), width=Inches(2.85))
    radius = 1.13
    cy = 2.09
    add_polar_descriptor(slide, 4.25, cy, radius, descriptor.sc_max, descriptor.sc_valid, "Conventional SC", descriptor.hidden_mask)
    divider = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.70), Inches(0.48), Inches(5.70), Inches(3.58))
    divider.line.color.rgb = common.rgb("#D4D8DC")
    divider.line.width = Pt(0.7)
    add_polar_descriptor(slide, 7.08, cy, radius, descriptor.upper_min, descriptor.upper_valid, "Up: min above split")
    add_polar_descriptor(slide, 9.90, cy, radius, descriptor.lower_max, descriptor.lower_valid, "Down: max below split")
    add_colorbar(slide, 11.62, 0.82, 0.14, 2.55)
    prs.save(ROOF_OUT)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument("--query-dir", type=Path, required=True)
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    common.build()
    build_principle()
    build_roof(args.query_dir)
    print(f"Wrote {PRINCIPLE_OUT}")
    print(f"Wrote {ROOF_OUT}")


if __name__ == "__main__":
    main()
