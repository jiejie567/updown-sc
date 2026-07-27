#!/usr/bin/env python3
"""Create an editable PowerPoint version of the UpDown-SC pipeline.

Every visible element is a native PowerPoint shape or text box.  No raster or
vector figure is embedded, so labels, colors, arrows, points, and panels can be
edited independently in PowerPoint or WPS Office.
"""

from __future__ import annotations

from math import cos, pi, sin
from pathlib import Path

import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parent / "updown_sc_pipeline_editable.pptx"


def rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


BLUE = "#0F4D92"
BLUE_2 = "#3775BA"
BLUE_LIGHT = "#DCE8F5"
TEAL = "#2A8C8C"
TEAL_LIGHT = "#D9EFEA"
RED = "#B64342"
ORANGE = "#D97706"
PURPLE = "#7A5AA6"
PURPLE_LIGHT = "#E9E1F2"
INK = "#25313C"
MID = "#65717C"
LIGHT = "#EEF2F5"
BORDER = "#C9D1D8"
WHITE = "#FFFFFF"
PANEL = "#FAFBFC"
HERO = "#F7FAFD"


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: float = 8,
    color: str = INK,
    bold: bool = False,
    align=PP_ALIGN.CENTER,
    font: str = "Arial",
    margin: float = 0.01,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str = WHITE,
    line: str = BORDER,
    line_width: float = 0.8,
    radius: bool = False,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(line_width)
    if radius:
        shape.adjustments[0] = 0.12
    return shape


def add_oval(slide, x: float, y: float, w: float, h: float, fill: str | None, line: str, width: float = 0.7):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(width)
    return shape


def add_line(
    slide,
    x0: float,
    y0: float,
    x1: float,
    y1: float,
    color: str = MID,
    width: float = 1.0,
    dashed: bool = False,
):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x0),
        Inches(y0),
        Inches(x1),
        Inches(y1),
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dashed:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return line


def add_right_arrow(slide, x: float, y: float, w: float, h: float, color: str = BLUE, rotation: float = 0):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = rgb(color)
    arrow.line.fill.background()
    arrow.rotation = rotation
    arrow.adjustments[0] = 0.55
    arrow.adjustments[1] = 0.72
    return arrow


def add_tag(slide, x: float, y: float, w: float, h: float, text: str, fill: str, color: str = INK, size: float = 7.5):
    add_rect(slide, x, y, w, h, fill=fill, line=fill, line_width=0, radius=True)
    return add_text(slide, x, y, w, h, text, size=size, color=color)


def add_panel(slide, x: float, y: float, w: float, h: float, label: str, title: str, fill: str = PANEL):
    add_rect(slide, x, y, w, h, fill=fill, line=BORDER, line_width=0.9)
    add_text(slide, x + 0.07, y + 0.05, 0.20, 0.25, label, size=11, bold=True, align=PP_ALIGN.LEFT)
    add_text(slide, x + 0.29, y + 0.05, w - 0.34, 0.25, title, size=9.3, bold=True, align=PP_ALIGN.LEFT)


def add_cloud(slide, x: float, y: float, w: float, h: float, tilt_deg: float = 0, density: float = 1.0):
    wall_x = np.linspace(-0.44, 0.44, max(7, int(12 * density)))
    floor = np.column_stack([wall_x, -0.31 + 0.035 * np.cos(8 * wall_x)])
    ceiling = np.column_stack([wall_x, 0.32 + 0.025 * np.sin(7 * wall_x)])
    shelf_y = np.linspace(-0.22, 0.08, max(4, int(7 * density)))
    shelf = np.column_stack([np.full_like(shelf_y, 0.25), shelf_y])
    doorway_y = np.linspace(-0.30, 0.22, max(5, int(9 * density)))
    doorway = np.column_stack([np.full_like(doorway_y, -0.33), doorway_y])
    pts = np.vstack([floor, ceiling, shelf, doorway])
    a = np.deg2rad(tilt_deg)
    rot = np.array([[np.cos(a), -np.sin(a)], [np.sin(a), np.cos(a)]])
    pts = pts @ rot.T
    for px, py in pts:
        cx = x + w * (0.5 + px)
        cy = y + h * (0.5 - py)
        add_oval(slide, cx - 0.025, cy - 0.025, 0.05, 0.05, BLUE_2, BLUE_2, 0)
    add_oval(slide, x + w * 0.5 - 0.025, y + h * 0.5 - 0.025, 0.05, 0.05, ORANGE, ORANGE, 0)


def add_axes(slide, ox: float, oy: float, angle_deg: float, scale: float = 0.28):
    a = angle_deg * pi / 180
    x1, y1 = ox + scale * cos(a), oy - scale * sin(a)
    z1, z1y = ox - scale * sin(a), oy - scale * cos(a)
    add_line(slide, ox, oy, x1, y1, MID, 1.1)
    add_line(slide, ox, oy, z1, z1y, ORANGE, 1.4)
    add_text(slide, x1 - 0.02, y1 - 0.05, 0.18, 0.16, "x", 7, MID)
    add_text(slide, z1 - 0.08, z1y - 0.10, 0.18, 0.16, "z", 7, ORANGE)


def add_descriptor_strip(slide, x: float, y: float, w: float, h: float, family: str, pattern: tuple[int, ...]):
    blue = ("#E7EFF8", "#BFD3EA", "#80A8D1", BLUE_2)
    teal = ("#E8F4F1", "#BEDFD8", "#79B9AE", TEAL)
    red = ("#F9E8E6", "#EBC3BF", "#D77F78", RED)
    palette = blue if family == "blue" else red if family == "red" else teal
    cw = w / len(pattern)
    for i, level in enumerate(pattern):
        add_rect(slide, x + i * cw, y, cw - 0.008, h, fill=palette[level], line=WHITE, line_width=0.2)


def add_polar_grid(slide, x: float, y: float, size: float):
    for frac in (1.0, 0.68, 0.36):
        inset = size * (1 - frac) / 2
        add_oval(slide, x + inset, y + inset, size * frac, size * frac, None, BLUE_2, 0.8)
    cx, cy = x + size / 2, y + size / 2
    for deg in range(0, 360, 45):
        a = deg * pi / 180
        add_line(slide, cx, cy, cx + size / 2 * cos(a), cy + size / 2 * sin(a), BLUE_2, 0.55)
    add_oval(slide, cx - 0.045, cy - 0.045, 0.09, 0.09, ORANGE, ORANGE, 0)


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(3.55)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = rgb(WHITE)

    py, ph = 0.08, 3.36
    xs = [0.06, 2.77, 6.96, 10.18]
    ws = [2.45, 3.93, 2.96, 3.09]
    for x, w, label, title, fill in zip(
        xs,
        ws,
        "abcd",
        ("Gravity-canonicalized scan", "Adaptive dual-envelope descriptor", "Masked place–yaw retrieval", "Optional 3-D verification"),
        (PANEL, HERO, PANEL, PANEL),
    ):
        add_panel(slide, x, py, w, ph, label, title, fill)

    # a) Common single-frame front end.
    x = xs[0]
    add_text(slide, x + 0.18, 0.56, 0.72, 0.20, "IMU up", 7.2, ORANGE)
    add_right_arrow(slide, x + 0.39, 0.76, 0.28, 0.11, ORANGE, 55)
    add_cloud(slide, x + 0.06, 0.91, 0.94, 0.84, -17, 0.62)
    add_axes(slide, x + 0.63, 1.56, -17, 0.25)
    add_right_arrow(slide, x + 1.00, 1.25, 0.34, 0.13, BLUE)
    add_text(slide, x + 0.91, 0.95, 0.55, 0.22, "deskew + Rg", 7.4, BLUE, True)
    add_cloud(slide, x + 1.35, 0.91, 0.94, 0.84, 0, 0.62)
    add_axes(slide, x + 1.92, 1.56, 0, 0.25)
    add_tag(slide, x + 0.15, 2.15, 0.86, 0.27, "blind 0.3 m", LIGHT)
    add_tag(slide, x + 1.35, 2.15, 0.88, 0.27, "body frame", BLUE_LIGHT, BLUE)
    add_text(slide, x + 0.33, 2.68, 1.79, 0.22, "one 0.1 s scan", 8.1, INK)
    add_text(slide, x + 0.20, 2.96, 2.05, 0.22, "roll/pitch removed · yaw free", 7.2, MID)

    # b) Adaptive dual-envelope descriptor (hero).
    x = xs[1]
    add_polar_grid(slide, x + 0.28, 0.91, 1.48)
    add_text(slide, x + 0.43, 2.42, 1.18, 0.22, "polar cells", 7.2, MID)
    add_line(slide, x + 1.93, 1.72, x + 3.64, 1.72, ORANGE, 1.2, True)
    add_text(slide, x + 2.35, 1.48, 0.94, 0.20, "cell-balanced τg", 7.0, ORANGE)
    add_line(slide, x + 2.00, 1.98, x + 3.30, 1.98, BLUE, 3)
    for xx, hh, shade in ((2.18, 0.58, BLUE_2), (2.57, 0.38, "#80A8D1"), (2.96, 0.50, "#5E91C6")):
        add_rect(slide, x + xx, 1.99, 0.20, hh, fill=shade, line=shade, line_width=0)
    add_text(slide, x + 2.01, 2.48, 1.28, 0.20, "Down: max below", 7.3, BLUE, True)
    add_line(slide, x + 2.00, 1.16, x + 3.30, 1.16, RED, 3)
    for xx, hh, shade in ((2.18, 0.42, "#D77F78"), (2.57, 0.30, RED), (2.96, 0.49, "#CB625C")):
        add_rect(slide, x + xx, 1.16 - hh, 0.20, hh, fill=shade, line=shade, line_width=0)
    add_text(slide, x + 2.01, 0.84, 1.28, 0.20, "Up: min above", 7.3, RED, True)
    # PowerPoint y grows downward: overhead/red is the upper row and
    # lower-middle/blue is the lower row.
    add_text(slide, x + 0.22, 2.73, 0.56, 0.18, "Up", 6.8, RED, True, PP_ALIGN.RIGHT)
    add_descriptor_strip(slide, x + 0.84, 2.74, 2.58, 0.18, "red", (3, 2, 0, 1, 3, 2, 2, 1))
    add_text(slide, x + 0.22, 2.98, 0.56, 0.18, "Down", 6.8, BLUE, True, PP_ALIGN.RIGHT)
    add_descriptor_strip(slide, x + 0.84, 2.99, 2.58, 0.18, "blue", (1, 3, 2, 0, 2, 3, 1, 2))
    add_text(slide, x + 0.57, 3.20, 2.80, 0.18, "overhead cannot overwrite lower geometry", 7.0, INK, True)

    # c) Masked SC-style place/yaw retrieval.
    x = xs[2]
    add_text(slide, x + 0.15, 0.61, 1.35, 0.20, "ring-key Top-K", 7.6, MID, align=PP_ALIGN.LEFT)
    for i, yy in enumerate((0.92, 1.37, 1.82)):
        selected = i == 1
        add_rect(slide, x + 0.20, yy, 1.05, 0.31, fill=BLUE_LIGHT if selected else WHITE, line=BLUE if selected else BORDER, line_width=1.0)
        add_descriptor_strip(slide, x + 0.34, yy + 0.10, 0.77, 0.10, "blue", (1 + i % 2, 3, 0, 2, 1, 2))
    add_right_arrow(slide, x + 1.34, 1.43, 0.35, 0.14, BLUE)
    add_oval(slide, x + 1.80, 1.15, 0.72, 0.72, None, PURPLE, 1.6)
    add_right_arrow(slide, x + 2.18, 1.13, 0.24, 0.11, PURPLE, 42)
    add_text(slide, x + 1.86, 1.40, 0.58, 0.20, "±3", 8.2, PURPLE, True)
    add_text(slide, x + 0.41, 2.36, 2.10, 0.21, "joint-valid rings", 8.0, INK, True)
    add_text(slide, x + 0.30, 2.62, 2.32, 0.22, "mask support × height cosine", 7.2, MID)
    add_tag(slide, x + 0.30, 2.93, 2.32, 0.27, "weighted two-channel distance", PURPLE_LIGHT, PURPLE, 7.2)

    # d) Six-degree-of-freedom geometric verification.
    x = xs[3]
    add_text(slide, x + 0.55, 0.65, 2.00, 0.22, "place + yaw seed", 8.2, PURPLE, True)
    add_rect(slide, x + 0.50, 1.10, 1.12, 0.84, fill=BLUE_LIGHT, line=BLUE, line_width=1.2)
    add_rect(slide, x + 1.40, 1.35, 1.12, 0.84, fill=WHITE, line=ORANGE, line_width=1.6)
    add_right_arrow(slide, x + 1.07, 0.91, 0.52, 0.14, PURPLE)
    add_text(slide, x + 0.61, 2.39, 1.90, 0.22, "point-to-map ICP", 8.2, INK, True)
    add_tag(slide, x + 0.82, 2.68, 1.48, 0.28, "accept / reject", TEAL_LIGHT, TEAL, 8.0)
    add_text(slide, x + 0.38, 3.04, 2.34, 0.20, "loop closure · relocalization", 6.8, MID)

    # Editable inter-panel arrows.
    for i in range(3):
        start = xs[i] + ws[i] + 0.08
        end = xs[i + 1] - 0.08
        add_right_arrow(slide, start, 1.68, end - start, 0.16, BLUE)

    prs.save(OUT)
    print(f"Wrote {OUT}")


def build_place_recognition_pipeline() -> None:
    """Create the editable loop-closure/relocalization front end used by the manuscript."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(3.02)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = rgb(WHITE)

    xs = [0.08 + 2.19 * index for index in range(6)]
    width, y, height = 1.82, 0.12, 2.76
    titles = (
        "LIO preprocessing",
        "Gravity canonicalization",
        "Dual-envelope SCD",
        "Place + yaw retrieval",
        "6-DoF ICP verification",
        "Successful place recognition",
    )
    fills = ("#F7FAFC", "#FFF9F1", "#F1F8F8", "#F7F4FB", "#FCF8F2", "#F3FAF6")
    badges = (BLUE_2, ORANGE, TEAL, PURPLE, ORANGE, "#2E8B57")
    for index, (x, title, fill, badge) in enumerate(zip(xs, titles, fills, badges), start=1):
        add_rect(slide, x + 0.025, y + 0.030, width, height, fill="#D9E0E6", line="#D9E0E6", line_width=0, radius=True)
        add_rect(slide, x, y, width, height, fill=fill, line="#C5CFD8", line_width=0.9, radius=True)
        add_line(slide, x + 0.10, y + 0.43, x + width - 0.10, y + 0.43, badge, 1.4)
        add_oval(slide, x + 0.08, y + 0.08, 0.25, 0.25, WHITE, badge, 1.2)
        add_text(slide, x + 0.08, y + 0.08, 0.25, 0.25, str(index), 8.0, badge, True)
        add_text(slide, x + 0.38, y + 0.04, 1.34, 0.36, title, 8.4, INK, True)

    # 1) Deskewed cloud and gravity from the LIO front end.
    x = xs[0]
    add_cloud(slide, x + 0.26, 0.82, 1.20, 0.86, -8, 0.70)
    add_right_arrow(slide, x + 0.21, 0.68, 0.32, 0.12, ORANGE, 90)
    add_text(slide, x + 0.10, 0.49, 0.54, 0.18, "gravity", 6.8, ORANGE, True)
    add_text(slide, x + 0.28, 1.82, 1.26, 0.22, "deskewed scan", 8.0, INK, True)
    add_text(slide, x + 0.22, 2.43, 1.38, 0.20, "LIO + IMU gravity", 7.0, MID)

    # 2) Gravity canonicalization.
    x = xs[1]
    add_axes(slide, x + 0.48, 1.38, -18, 0.34)
    add_right_arrow(slide, x + 0.76, 1.28, 0.30, 0.13, BLUE)
    add_text(slide, x + 0.71, 1.03, 0.42, 0.19, "align", 6.8, BLUE, True)
    add_axes(slide, x + 1.34, 1.38, 0, 0.34)
    add_text(slide, x + 0.18, 1.82, 1.48, 0.22, "gravity-aligned cloud", 7.8, INK, True)
    add_text(slide, x + 0.12, 2.43, 1.58, 0.20, "roll/pitch fixed · yaw free", 6.8, MID)

    # 3) Dual-envelope construction.
    x = xs[2]
    add_polar_grid(slide, x + 0.14, 0.79, 0.82)
    add_text(slide, x + 0.14, 1.66, 0.82, 0.18, "polar cells", 6.7, MID)
    # PowerPoint y grows downward: place overhead/red visually above
    # lower-middle/blue, matching the manuscript schematics.
    add_text(slide, x + 1.08, 0.68, 0.68, 0.18, "Up (overhead)", 6.0, RED, True)
    add_descriptor_strip(slide, x + 1.17, 0.91, 0.50, 0.13, "red", (3, 2, 0, 1, 3, 2))
    add_text(slide, x + 1.04, 1.45, 0.76, 0.18, "Down (lower)", 5.8, BLUE, True)
    add_descriptor_strip(slide, x + 1.17, 1.28, 0.50, 0.13, "blue", (1, 3, 2, 0, 2, 3))
    add_text(slide, x + 0.16, 1.98, 1.50, 0.20, "adaptive split + masks", 7.6, INK, True)
    add_text(slide, x + 0.14, 2.43, 1.54, 0.20, "map/query same definition", 6.8, MID)

    # 4) Keyframe-database retrieval and yaw estimation.
    x = xs[3]
    add_text(slide, x + 0.20, 0.53, 0.86, 0.20, "keyframe DB", 7.0, MID)
    for index, yy in enumerate((0.84, 1.15, 1.46)):
        selected = index == 1
        add_rect(slide, x + 0.25, yy, 0.80, 0.24, BLUE_LIGHT if selected else WHITE, BLUE if selected else BORDER, 0.8)
        add_descriptor_strip(slide, x + 0.35, yy + 0.045, 0.60, 0.055, "red", (2, 3, 0, 1, 2, 1))
        add_descriptor_strip(slide, x + 0.35, yy + 0.140, 0.60, 0.055, "blue", (1 + index % 2, 3, 0, 2, 1, 2))
    add_right_arrow(slide, x + 1.08, 1.21, 0.22, 0.11, BLUE)
    add_oval(slide, x + 1.34, 1.02, 0.38, 0.38, None, PURPLE, 1.2)
    add_text(slide, x + 1.35, 1.10, 0.36, 0.20, "yaw", 6.8, PURPLE, True)
    add_text(slide, x + 0.17, 1.98, 1.48, 0.20, "masked SCD matching", 7.6, INK, True)
    add_text(slide, x + 0.20, 2.43, 1.42, 0.20, "ranked places + yaw", 6.7, MID)

    # 5) ICP refinement and health checks.
    x = xs[4]
    add_rect(slide, x + 0.30, 0.91, 0.76, 0.61, BLUE_LIGHT, BLUE, 1.0)
    add_rect(slide, x + 0.70, 1.13, 0.76, 0.61, WHITE, ORANGE, 1.2)
    add_right_arrow(slide, x + 0.69, 0.70, 0.38, 0.12, PURPLE)
    add_oval(slide, x + 1.35, 0.62, 0.28, 0.28, TEAL_LIGHT, "#2E8B57", 0.9)
    add_text(slide, x + 1.35, 0.63, 0.28, 0.24, "✓", 9.0, "#2E8B57", True)
    add_text(slide, x + 0.25, 1.98, 1.32, 0.20, "point-to-map ICP", 7.8, INK, True)
    add_text(slide, x + 0.20, 2.38, 1.42, 0.38, "overlap · fitness\nconvergence", 6.6, MID)

    # 6) Accepted map-frame pose.
    x = xs[5]
    path = ((0.27, 1.64), (0.50, 1.24), (0.80, 1.48), (1.03, 0.93), (1.36, 1.24), (1.58, 0.70))
    for first, second in zip(path[:-1], path[1:]):
        add_line(slide, x + first[0], first[1], x + second[0], second[1], BORDER, 3.0)
    for px, py in path[:-1]:
        add_oval(slide, x + px - 0.05, py - 0.05, 0.10, 0.10, WHITE, MID, 0.8)
    pose = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x + 1.20), Inches(0.92), Inches(0.30), Inches(0.24))
    pose.rotation = 90
    pose.fill.solid()
    pose.fill.fore_color.rgb = rgb("#2E8B57")
    pose.line.fill.background()
    add_oval(slide, x + 1.38, 1.55, 0.30, 0.30, TEAL_LIGHT, "#2E8B57", 0.9)
    add_text(slide, x + 1.38, 1.58, 0.30, 0.24, "✓", 9.0, "#2E8B57", True)
    add_text(slide, x + 0.14, 1.98, 1.54, 0.20, "recognized map place", 7.5, INK, True)
    add_text(slide, x + 0.18, 2.43, 1.46, 0.20, "recognition succeeded", 6.8, MID)

    for index in range(5):
        start = xs[index] + width + 0.03
        end = xs[index + 1] - 0.03
        add_right_arrow(slide, start, 1.39, end - start, 0.13, "#2E8B57" if index == 4 else BLUE)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build_place_recognition_pipeline()
