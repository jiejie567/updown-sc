#!/usr/bin/env python3
"""Generate the page-one overview and aggregate retrieval figure.

Figure contract
---------------
Core conclusion:
    Dual envelopes prevent ceiling maxima from erasing lower indoor structure
    and improve first-choice place retrieval under a shared gravity front end.
Figure archetype:
    Fig. 1 is a schematic-led overview; the retrieval figure is an asymmetric
    mixed-modality figure with spatial error distributions and rank curves.
Target/output:
    ICRA/IEEE two-column paper, editable SVG/PDF and native-shape PowerPoint.
Backend:
    Python/Matplotlib.
Final size:
    Fig. 1 is 89 mm wide; retrieval subfigures compose to 183 mm wide.
Panel map:
    Overview: the same indoor returns -> ceiling-dominated SC versus
    complementary UpDown-SC envelopes.
    Retrieval a/b: every eligible IH Top-1 result for SC+G / UpDown-SC+G.
    Retrieval c: Recall@K, K=1..5, on IH and CH.
Evidence hierarchy:
    Hero evidence is the all-query spatial distribution and Recall@1..5 curves;
    the overview defines the method and its pose-seed output.
Statistics needed:
    Descriptive recall only; n=320 IH queries and n=148 CH queries.
Source data needed:
    Audited per-query CSVs already used by the manuscript table.
Image-integrity notes:
    No raster data are modified. Every eligible query contributes once.
Reviewer risk:
    Rank curves must reproduce the manuscript's R@1 and R@5 table entries, and
    the conceptual overview must not be mistaken for measured performance.
"""

from __future__ import annotations

import csv
import json
import math
from pathlib import Path

import matplotlib as mpl
import matplotlib.pyplot as plt
import numpy as np
from matplotlib.lines import Line2D
from matplotlib.patches import Arc, Circle, FancyArrowPatch, Rectangle
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

import generate_roof_contamination_figure as roof


OUT = Path(__file__).resolve().parent
EXP = Path("${UPDOWN_SC_ROOT}/icra2027_runtime/experiments")

IH_MULTI = EXP / "gravity_transfer_2m_20260718/ih/results/per_query.csv"
IH_IRIS = EXP / "gravity_transfer_2m_20260718/ih/results/lidar_iris_per_query.csv"
IH_RING = EXP / "gravity_transfer_2m_20260718/ih/ringpp/results/ringpp_per_query.csv"
IH_OURS = EXP / "updown_weight_ablation_real_20260721/selected/ih_gravity/paper_per_query.csv"
IH_MAP = EXP / "private_two_bag_2m/protocol/map/selected_poses_tum.txt"
HERO_QUERY_DIR = EXP / "baseline_20260715/queries/loc_2_floor_continuous"
HERO_QUERY_ID = 112
HERO_SPLIT_HEIGHT_M = 2.1
HERO_RANGE_M = 8.0
HERO_OVERHEAD_RANGE_M = 4.0
HERO_ELEV_DEG = 16.0
HERO_SIDE_AZIM_DEG = 90.0
HERO_AZIM_DEG = -45.0
HERO_XY_PERCENTILES = (2.0, 98.0)
HERO_COLOR_MIN_M = -0.5
HERO_CEILING_SEARCH_MAX_M = 4.5
HERO_CEILING_BIN_M = 0.02
HERO_CEILING_MARGIN_M = 0.12

CH_MULTI = EXP / "gravity_transfer_2m_20260718/ch/results/per_query.csv"
CH_IRIS = EXP / "gravity_transfer_2m_20260718/ch/results/lidar_iris_per_query.csv"
CH_RING = EXP / "gravity_transfer_2m_20260718/ch/ringpp/results/ringpp_per_query.csv"
CH_OURS = EXP / "updown_weight_ablation_real_20260721/selected/ch_gravity/paper_per_query.csv"

INK = "#24313D"
MID = "#6B7783"
LIGHT = "#D6DEE5"
PALE = "#EEF3F7"
BLUE = "#1764AA"
BLUE_LIGHT = "#DDEAF6"
TEAL = "#238D8E"
TEAL_LIGHT = "#D7EEEA"
ORANGE = "#D97706"
PURPLE = "#7653A6"
RED = "#C94945"
WHITE = "#FFFFFF"
UPPER = RED

HEIGHT_CMAP = mpl.colors.LinearSegmentedColormap.from_list(
    "height",
    ["#263B73", "#3775BA", "#D9EFEA", "#F4C66A", "#B64342"],
)

METHODS = {
    "SC": "#7A8792",
    "SC++": "#4F9A9A",
    "LiDAR Iris": "#D58A24",
    "RING++": "#7653A6",
    "UpDown-SC": BLUE,
}

DATASETS = {
    "IH": {
        "n": 320,
        "sources": {
            "SC": (IH_MULTI, "SC + G"),
            "SC++": (IH_MULTI, "SC++ (PC) + G"),
            "LiDAR Iris": (IH_IRIS, "LiDAR-Iris + G"),
            "RING++": (IH_RING, "RING++ + G"),
            "UpDown-SC": (IH_OURS, "UpDown-SC pure adaptive"),
        },
    },
    "CH": {
        "n": 148,
        "sources": {
            "SC": (CH_MULTI, "SC + G"),
            "SC++": (CH_MULTI, "SC++ (PC) + G"),
            "LiDAR Iris": (CH_IRIS, "LiDAR-Iris + G"),
            "RING++": (CH_RING, "RING++ + G"),
            "UpDown-SC": (CH_OURS, "UpDown-SC pure adaptive"),
        },
    },
}

mpl.rcParams.update(
    {
        "font.family": "sans-serif",
        "font.sans-serif": ["Arial", "Helvetica", "DejaVu Sans", "sans-serif"],
        "font.size": 7.0,
        "axes.labelsize": 7.0,
        "axes.titlesize": 7.5,
        "xtick.labelsize": 6.4,
        "ytick.labelsize": 6.4,
        "legend.fontsize": 6.3,
        "axes.linewidth": 0.7,
        "axes.spines.top": False,
        "axes.spines.right": False,
        "svg.fonttype": "none",
        "pdf.fonttype": 42,
    }
)


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.lstrip("#")
    return RGBColor(*[int(value[i : i + 2], 16) for i in (0, 2, 4)])


def read_rows(path: Path, algorithm: str) -> list[dict[str, str]]:
    with path.open(newline="") as handle:
        rows = [row for row in csv.DictReader(handle) if row["algorithm"] == algorithm]
    if not rows:
        raise RuntimeError(f"No rows for {algorithm!r} in {path}")
    return rows


def rank_value(row: dict[str, str]) -> int:
    value = row.get("first_correct_rank", "").strip()
    if not value or value.lower() in {"nan", "none", "inf"}:
        return 10**9
    return int(float(value))


def recall_curve(rows: list[dict[str, str]]) -> np.ndarray:
    ranks = np.array([rank_value(row) for row in rows], dtype=np.int64)
    return np.array([100.0 * np.mean(ranks <= k) for k in range(1, 6)])


def load_all() -> dict[str, dict[str, list[dict[str, str]]]]:
    loaded: dict[str, dict[str, list[dict[str, str]]]] = {}
    for dataset, spec in DATASETS.items():
        loaded[dataset] = {}
        for method, (path, algorithm) in spec["sources"].items():
            rows = read_rows(path, algorithm)
            if len(rows) != spec["n"]:
                raise RuntimeError(
                    f"{dataset}/{method}: expected {spec['n']} rows, got {len(rows)}"
                )
            loaded[dataset][method] = rows
    return loaded


def save(fig: mpl.figure.Figure, stem: str, dpi: int = 600) -> None:
    # Keep any intentionally rasterized artists at publication resolution.
    # The measured Fig. 1 cloud itself remains vector in PDF/SVG.
    fig.savefig(OUT / f"{stem}.svg", dpi=dpi, bbox_inches="tight", pad_inches=0.02)
    fig.savefig(OUT / f"{stem}.pdf", dpi=dpi, bbox_inches="tight", pad_inches=0.02)
    fig.savefig(OUT / f"{stem}.png", dpi=dpi, bbox_inches="tight", pad_inches=0.02)
    fig.savefig(OUT / f"{stem}.tiff", dpi=dpi, bbox_inches="tight", pad_inches=0.02)
    plt.close(fig)


def arrow(
    ax: mpl.axes.Axes,
    xy0: tuple[float, float],
    xy1: tuple[float, float],
    color: str = BLUE,
    lw: float = 1.1,
    scale: float = 8.0,
) -> None:
    ax.add_patch(
        FancyArrowPatch(
            xy0,
            xy1,
            arrowstyle="-|>",
            mutation_scale=scale,
            color=color,
            linewidth=lw,
            shrinkA=0,
            shrinkB=0,
        )
    )


def descriptor_matrix(
    ax: mpl.axes.Axes,
    x: float,
    y: float,
    width: float,
    height: float,
    color: str,
    shift: int = 0,
) -> None:
    base = np.array(
        [
            [0, 1, 3, 2, 0, 1, 2, 3, 1, 0],
            [1, 2, 3, 1, 0, 2, 3, 2, 1, 0],
            [0, 1, 2, 0, 1, 3, 2, 1, 0, 2],
            [1, 0, 2, 1, 2, 3, 1, 0, 2, 1],
        ],
        dtype=int,
    )
    base = np.roll(base, shift, axis=1)
    rows, cols = base.shape
    for i in range(rows):
        for j in range(cols):
            alpha = (0.08, 0.30, 0.58, 0.90)[base[i, j]]
            ax.add_patch(
                Rectangle(
                    (x + j * width / cols, y + (rows - i - 1) * height / rows),
                    width / cols,
                    height / rows,
                    facecolor=color,
                    edgecolor=WHITE,
                    linewidth=0.25,
                    alpha=alpha,
                )
            )
    ax.add_patch(
        Rectangle((x, y), width, height, fill=False, edgecolor=LIGHT, linewidth=0.6)
    )


def load_hero_cloud() -> tuple[roof.QueryRecord, np.ndarray]:
    """Load the same measured query used by the descriptor-detail figures."""
    roof.SPLIT_HEIGHT_M = HERO_SPLIT_HEIGHT_M
    record, descriptor, _ = roof.select_query(
        roof.load_queries(HERO_QUERY_DIR), HERO_QUERY_ID
    )
    return record, descriptor.render_xyz


def ceiling_display_filter(xyz: np.ndarray) -> tuple[np.ndarray, float, float]:
    """Retain the dominant ceiling plane and remove only returns above it."""
    height = xyz[:, 2]
    candidates = height[
        (height >= HERO_SPLIT_HEIGHT_M + 0.2)
        & (height <= HERO_CEILING_SEARCH_MAX_M)
    ]
    bins = np.arange(
        HERO_SPLIT_HEIGHT_M + 0.2,
        HERO_CEILING_SEARCH_MAX_M + HERO_CEILING_BIN_M,
        HERO_CEILING_BIN_M,
    )
    counts, edges = np.histogram(candidates, bins=bins)
    peak = int(np.argmax(counts))
    ceiling_height = 0.5 * (edges[peak] + edges[peak + 1])
    ceiling_cutoff = ceiling_height + HERO_CEILING_MARGIN_M
    return xyz[height <= ceiling_cutoff], float(ceiling_height), float(ceiling_cutoff)


def prepare_hero_display(
    xyz: np.ndarray,
) -> tuple[np.ndarray, float, float, np.ndarray, np.ndarray, dict[str, float | int]]:
    """Apply the documented display-only crop for the measured hero cloud."""
    finite = np.isfinite(xyz).all(axis=1)
    xyz = xyz[finite]
    points_total = len(xyz)
    xyz, ceiling_height, ceiling_cutoff = ceiling_display_filter(xyz)
    points_after_ceiling = len(xyz)
    radius = np.hypot(xyz[:, 0], xyz[:, 1])
    xyz = xyz[radius <= HERO_RANGE_M]
    points_after_range = len(xyz)
    xy_low = np.percentile(xyz[:, :2], HERO_XY_PERCENTILES[0], axis=0)
    xy_high = np.percentile(xyz[:, :2], HERO_XY_PERCENTILES[1], axis=0)
    inside_xy = (
        (xyz[:, 0] >= xy_low[0])
        & (xyz[:, 0] <= xy_high[0])
        & (xyz[:, 1] >= xy_low[1])
        & (xyz[:, 1] <= xy_high[1])
    )
    xyz = xyz[inside_xy]
    stats = {
        "points_total": int(points_total),
        "dominant_ceiling_height_m": ceiling_height,
        "ceiling_display_cutoff_m": ceiling_cutoff,
        "points_after_ceiling_filter": int(points_after_ceiling),
        "points_removed_above_ceiling": int(points_total - points_after_ceiling),
        "radial_display_crop_m": HERO_RANGE_M,
        "points_after_radial_crop": int(points_after_range),
        "xy_display_percentiles": list(HERO_XY_PERCENTILES),
        "points_displayed": int(len(xyz)),
        "display_x_low_m": float(xy_low[0]),
        "display_x_high_m": float(xy_high[0]),
        "display_y_low_m": float(xy_low[1]),
        "display_y_high_m": float(xy_high[1]),
        "display_z_low_m": -0.25,
        "display_z_high_m": ceiling_cutoff + 0.04,
    }
    return xyz, ceiling_height, ceiling_cutoff, xy_low, xy_high, stats


def orthographic_project(xyz: np.ndarray) -> tuple[np.ndarray, np.ndarray, np.ndarray]:
    """Project XYZ to the selected oblique view without a perspective camera."""
    azimuth = np.deg2rad(HERO_AZIM_DEG)
    elevation = np.deg2rad(HERO_ELEV_DEG)
    horizontal = -np.sin(azimuth) * xyz[:, 0] + np.cos(azimuth) * xyz[:, 1]
    vertical = (
        -np.cos(azimuth) * np.sin(elevation) * xyz[:, 0]
        - np.sin(azimuth) * np.sin(elevation) * xyz[:, 1]
        + np.cos(elevation) * xyz[:, 2]
    )
    depth = (
        np.cos(azimuth) * np.cos(elevation) * xyz[:, 0]
        + np.sin(azimuth) * np.cos(elevation) * xyz[:, 1]
        + np.sin(elevation) * xyz[:, 2]
    )
    return horizontal, vertical, depth


def draw_measured_sideview(parent_ax, xyz: np.ndarray) -> dict[str, float | int]:
    """Draw the selected measured scan as a side elevation view."""
    xyz, _, ceiling_cutoff, xy_low, xy_high, stats = prepare_hero_display(xyz)
    stats["side_view_azimuth_deg"] = HERO_SIDE_AZIM_DEG
    height_norm = mpl.colors.Normalize(
        vmin=HERO_COLOR_MIN_M, vmax=ceiling_cutoff, clip=True
    )

    # Reserve a clean title band above the measured cloud at final column size.
    side_left, side_bottom, side_width, side_height = 0.025, 0.610, 0.835, 0.315
    side = parent_ax.inset_axes([side_left, side_bottom, side_width, side_height])
    azimuth = np.deg2rad(HERO_SIDE_AZIM_DEG)
    horizontal = xyz[:, 0] * np.cos(azimuth) + xyz[:, 1] * np.sin(azimuth)
    depth = -xyz[:, 0] * np.sin(azimuth) + xyz[:, 1] * np.cos(azimuth)
    vertical = xyz[:, 2]
    order = np.argsort(depth)[::-1]
    side.scatter(
        horizontal[order],
        vertical[order],
        s=0.76,
        c=vertical[order],
        cmap=HEIGHT_CMAP,
        norm=height_norm,
        linewidths=0,
        rasterized=False,
    )
    h_pad = 0.02 * np.ptp(horizontal)
    h_low = float(horizontal.min() - h_pad)
    h_high = float(horizontal.max() + h_pad)
    side.plot(
        [h_low, h_high],
        [HERO_SPLIT_HEIGHT_M] * 2,
        color=ORANGE,
        linewidth=0.7,
        linestyle=(0, (3, 2)),
    )
    side.set_xlim(h_low, h_high)
    side.set_ylim(-0.40, ceiling_cutoff + 0.10)
    side.set_aspect("equal", adjustable="box")
    side.axis("off")
    side.patch.set_alpha(0.0)
    parent_ax.text(0.800, 0.612, r"$r\leq 8$ m", color=MID, fontsize=5.0)

    # Channel labels anchored on measured structure.
    def anchor(mask: np.ndarray, h_pct: float, v_pct: float) -> tuple[float, float]:
        h_sel, v_sel = horizontal[mask], vertical[mask]
        target = np.array([np.percentile(h_sel, h_pct), np.percentile(v_sel, v_pct)])
        idx = int(np.argmin(np.hypot(h_sel - target[0], v_sel - target[1])))
        return float(h_sel[idx]), float(v_sel[idx])

    up_mask = vertical > HERO_SPLIT_HEIGHT_M
    down_mask = (vertical > 0.3) & (vertical < 1.6)
    h_down = horizontal[down_mask]
    box_left = float(np.percentile(h_down, 3.0))
    box_right = float(np.percentile(h_down, 97.0))
    side.add_patch(
        Rectangle(
            (box_left, -0.08),
            box_right - box_left,
            1.62,
            facecolor="none",
            edgecolor=RED,
            linewidth=0.7,
            linestyle=(0, (4, 2.5)),
            zorder=6,
        )
    )
    label_box = {"boxstyle": "square,pad=0.15", "facecolor": WHITE,
                 "edgecolor": "none", "alpha": 0.85}
    side.annotate("Up: ceiling", anchor(up_mask, 18.0, 60.0),
                  xytext=(0.015, 0.93), textcoords="axes fraction",
                  ha="left", va="top", fontsize=5.4, color=UPPER,
                  fontweight="bold", bbox=label_box,
                  arrowprops={"arrowstyle": "-", "color": UPPER,
                              "linewidth": 0.55, "shrinkB": 1.5})
    side.text(box_right, -0.16, "Down: lower structure (desks, chairs)",
              ha="right", va="top", fontsize=5.4, color=BLUE,
              fontweight="bold", bbox=label_box)

    # Keep the height legend outside the measured point layer.
    cax = parent_ax.inset_axes([0.865, side_bottom + 0.018, 0.018, side_height - 0.036])
    colorbar = parent_ax.figure.colorbar(
        mpl.cm.ScalarMappable(norm=height_norm, cmap=HEIGHT_CMAP), cax=cax
    )
    colorbar.set_ticks([0.0, HERO_SPLIT_HEIGHT_M, ceiling_cutoff])
    colorbar.set_ticklabels(["0", "2.1", f"{ceiling_cutoff:.1f}"])
    cax.text(
        -0.55,
        HERO_SPLIT_HEIGHT_M,
        r"$\tau_g$",
        transform=mpl.transforms.blended_transform_factory(
            cax.transAxes, cax.transData),
        color=ORANGE,
        fontsize=5.2,
        ha="right",
        va="center",
    )
    colorbar.ax.tick_params(labelsize=5.0, length=1.8, width=0.5, pad=1.2)
    colorbar.outline.set_linewidth(0.45)
    colorbar.set_label(r"height $z_g$ (m)", fontsize=5.0, labelpad=1.5)

    return stats


def save_sideview_asset(xyz: np.ndarray) -> None:
    """Export the dense measured layer for insertion into editable PowerPoint."""
    xyz, _, ceiling_cutoff, xy_low, xy_high, _ = prepare_hero_display(xyz)
    height_norm = mpl.colors.Normalize(
        vmin=HERO_COLOR_MIN_M, vmax=ceiling_cutoff, clip=True
    )
    # Keep the auxiliary PPT raster at the same 87.6 mm single-column width
    # as the manuscript figure so automated final-size QA sees one contract.
    fig = plt.figure(figsize=(3.45, 1.16), facecolor=WHITE)
    ax = fig.add_axes([0.015, 0.04, 0.85, 0.92])
    horizontal, vertical, depth = orthographic_project(xyz)
    order = np.argsort(depth)[::-1]
    ax.scatter(
        horizontal[order],
        vertical[order],
        s=0.68,
        c=xyz[order, 2],
        cmap=HEIGHT_CMAP,
        norm=height_norm,
        linewidths=0,
        rasterized=True,
    )
    split_x = [xy_low[0], xy_high[0], xy_high[0], xy_low[0], xy_low[0]]
    split_y = [xy_low[1], xy_low[1], xy_high[1], xy_high[1], xy_low[1]]
    split_xyz = np.column_stack((split_x, split_y, [HERO_SPLIT_HEIGHT_M] * 5))
    split_horizontal, split_vertical, _ = orthographic_project(split_xyz)
    ax.plot(split_horizontal, split_vertical, color=ORANGE, linewidth=0.9, linestyle=(0, (3, 2)))
    h_pad = 0.02 * np.ptp(horizontal)
    v_pad = 0.02 * np.ptp(vertical)
    ax.set_xlim(float(horizontal.min() - h_pad), float(horizontal.max() + h_pad))
    ax.set_ylim(float(vertical.min() - v_pad), float(vertical.max() + v_pad))
    ax.set_aspect("equal", adjustable="box")
    ax.axis("off")
    cax = fig.add_axes([0.895, 0.20, 0.018, 0.66])
    colorbar = fig.colorbar(
        mpl.cm.ScalarMappable(norm=height_norm, cmap=HEIGHT_CMAP), cax=cax
    )
    colorbar.set_ticks([0.0, HERO_SPLIT_HEIGHT_M, ceiling_cutoff])
    colorbar.set_ticklabels(["0", "2.1", f"{ceiling_cutoff:.1f}"])
    cax.text(
        -0.55,
        HERO_SPLIT_HEIGHT_M,
        r"$\tau_g$",
        transform=mpl.transforms.blended_transform_factory(
            cax.transAxes, cax.transData),
        color=ORANGE,
        fontsize=5.2,
        ha="right",
        va="center",
    )
    colorbar.ax.tick_params(labelsize=5.0, length=1.8, width=0.5, pad=1.2)
    colorbar.outline.set_linewidth(0.45)
    colorbar.set_label(r"height $z_g$ (m)", fontsize=5.0, labelpad=1.5)
    fig.text(0.055, 0.10, r"$r\leq 8$ m", ha="left", va="bottom", fontsize=5.0, color=MID)
    fig.savefig(OUT / "updown_sc_hero_sideview.png", dpi=900, facecolor=WHITE)
    plt.close(fig)


def write_hero_metadata(record: roof.QueryRecord, stats: dict[str, float | int]) -> None:
    metadata = {
        "source": "user-selected candidate H (query 112) from hero_oblique_candidates",
        "query_id": record.query_id,
        "single_frame_duration_s": 0.1,
        "gravity_canonicalized": True,
        "projection": (
            f"orthographic side elevation view; azimuth "
            f"{HERO_SIDE_AZIM_DEG:.0f} degrees; gravity-aligned z"
        ),
        "split_height_m": HERO_SPLIT_HEIGHT_M,
        "color_semantics": {
            "quantity": "gravity-aligned height z_g in metres",
            "low": "blue",
            "high": "red",
            "display_min_m": HERO_COLOR_MIN_M,
            "display_max": "dominant ceiling height plus 0.12 m",
            "split": "orange dashed line at the physical split height",
        },
        "display_crop": stats,
        "image_integrity": (
            "For the Fig. 1 side-view display only, returns beyond 8 m, above the dominant ceiling-height "
            "histogram mode plus 0.12 m, or outside the stated XY display percentiles are omitted. "
            "Descriptor construction and every reported experiment remain unchanged. "
            "No point is moved, interpolated, locally recolored, or manually removed."
        ),
    }
    (OUT / "updown_sc_hero_metadata.json").write_text(json.dumps(metadata, indent=2) + "\n")


def make_overview(record: roof.QueryRecord, xyz: np.ndarray) -> None:
    """Page-one problem/solution figure, following the visual logic of SC papers."""
    fig, ax = plt.subplots(figsize=(3.45, 2.56), facecolor=WHITE)
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis("off")

    # The user-selected measured indoor-start scan supplies the oblique evidence.
    ax.text(
        0.50,
        0.982,
        "Measured indoor scan (side view)",
        ha="center",
        va="top",
        fontsize=8.0,
        weight="bold",
        color=INK,
    )
    side_stats = draw_measured_sideview(ax, xyz)

    # Branch to the two representations.
    arrow(ax, (0.36, 0.60), (0.26, 0.535), RED, 1.1, 7)
    arrow(ax, (0.64, 0.60), (0.74, 0.535), BLUE, 1.1, 7)
    ax.add_patch(Rectangle((0.035, 0.055), 0.445, 0.475, facecolor="#FFF8F7", edgecolor="#E7C6C3"))
    ax.add_patch(Rectangle((0.520, 0.055), 0.445, 0.475, facecolor="#F5F9FD", edgecolor="#BED4E8"))
    ax.text(0.258, 0.493, "Conventional SC", ha="center", weight="bold", fontsize=7.6, color=INK)
    ax.text(0.742, 0.493, "UpDown-SC", ha="center", weight="bold", fontsize=7.6, color=INK)

    # Conventional SC: every cell takes the ceiling maximum, producing a
    # nearly saturated and weakly distinctive height image.
    sc_ceiling_bottoms = (0.408, 0.398, 0.414, 0.402, 0.411, 0.399, 0.409)
    sc_lower_tops = (0.315, 0.340, 0.300, 0.355, 0.323, 0.345, 0.305)
    for x, bottom, top in zip(
            np.linspace(0.10, 0.42, 7), sc_ceiling_bottoms, sc_lower_tops):
        ax.plot([x, x], [bottom, 0.448], color=UPPER, lw=5.2, alpha=0.72)
        ax.plot([x, x], [0.282, top], color=LIGHT, lw=5.2)
    ax.plot([0.085, 0.435], [0.448, 0.448], color=UPPER, lw=1.0,
            ls=(0, (3, 2)))
    sc_pattern = (3, 3, 3, 3, 3, 3, 3, 2, 3, 3)
    descriptor_matrix(ax, 0.082, 0.125, 0.35, 0.090, UPPER, shift=0)
    # Overlay uniform cells to make saturation visually unambiguous.
    for j, level in enumerate(sc_pattern):
        ax.add_patch(
            Rectangle(
                (0.082 + j * 0.035, 0.125),
                0.0345,
                0.090,
                facecolor=UPPER,
                edgecolor=WHITE,
                linewidth=0.25,
                alpha=(0.62, 0.82)[level == 3],
            )
        )
    ax.text(0.258, 0.062, "ceiling-dominated", ha="center", fontsize=5.9, color=RED, weight="bold")

    # UpDown-SC: lower/middle structure remains distinctive while the overhead
    # layer is preserved independently.
    y_split = 0.367
    ax.plot([0.565, 0.925], [y_split, y_split], color=ORANGE, lw=1.0, ls=(0, (3, 2)))
    ax.text(0.93, y_split, r"$\tau_g$", fontsize=6.2, color=ORANGE, va="center")
    lower_heights = [0.315, 0.340, 0.300, 0.355, 0.323, 0.345, 0.305]
    upper_bottoms = [0.412, 0.424, 0.406, 0.427, 0.409, 0.422, 0.408]
    bar_x = np.linspace(0.59, 0.90, 7)
    for x, y, u in zip(bar_x, lower_heights, upper_bottoms):
        ax.plot([x, x], [0.282, y], color=BLUE, lw=5.0, alpha=0.78)
        ax.plot([x, x], [u, 0.448], color=UPPER, lw=5.0, alpha=0.78)
    # Envelope polylines: min height above the split, max height below it.
    ax.plot(np.concatenate(([0.575], bar_x, [0.905])),
            np.concatenate(([upper_bottoms[0]], upper_bottoms, [upper_bottoms[-1]])),
            color=UPPER, lw=0.9, ls=(0, (3, 2)), zorder=5)
    ax.plot(np.concatenate(([0.575], bar_x, [0.905])),
            np.concatenate(([lower_heights[0]], lower_heights, [lower_heights[-1]])),
            color=BLUE, lw=0.9, ls=(0, (3, 2)), zorder=5)
    # Preserve the physical vertical and color order used by Figs. 3--4:
    # overhead/red is the
    # upper row, while lower-middle/blue is the lower row.
    descriptor_matrix(ax, 0.562, 0.158, 0.36, 0.066, UPPER, shift=0)
    descriptor_matrix(ax, 0.562, 0.098, 0.36, 0.052, BLUE, shift=0)
    ax.text(0.742, 0.062, "structure retained", ha="center", fontsize=5.9, color=BLUE, weight="bold")

    save(fig, "updown_sc_hero")
    save_sideview_asset(xyz)
    write_hero_metadata(record, side_stats)


def spatial_arrays(rows: list[dict[str, str]]) -> tuple[np.ndarray, np.ndarray, np.ndarray]:
    truth = np.array([[float(row["truth_x"]), float(row["truth_y"])] for row in rows])
    top1 = np.array([[float(row["top1_x"]), float(row["top1_y"])] for row in rows])
    correct = np.array([rank_value(row) == 1 for row in rows], dtype=bool)
    return truth, top1, correct


def common_bounds(map_xy: np.ndarray, truth: np.ndarray, top1: np.ndarray) -> tuple[float, float, float, float]:
    all_xy = np.vstack((map_xy, truth, top1))
    low = np.nanmin(all_xy, axis=0)
    high = np.nanmax(all_xy, axis=0)
    span = np.maximum(high - low, 1.0)
    low -= 0.06 * span
    high += 0.06 * span
    return low[0], high[0], low[1], high[1]


def make_spatial_panel(
    map_xy: np.ndarray,
    rows: list[dict[str, str]],
    bounds: tuple[float, float, float, float],
    stem: str,
) -> None:
    truth, top1, correct = spatial_arrays(rows)
    fig, ax = plt.subplots(figsize=(2.22, 2.38), facecolor=WHITE)
    ax.plot(map_xy[:, 0], map_xy[:, 1], color="#C6CFD6", lw=1.0, zorder=1)
    ax.plot(truth[:, 0], truth[:, 1], color="#7D8993", lw=0.55, alpha=0.75, zorder=2)
    wrong = ~correct
    for q, m in zip(truth[wrong], top1[wrong]):
        ax.plot([q[0], m[0]], [q[1], m[1]], color=RED, lw=0.35, alpha=0.18, zorder=0)
    ax.scatter(
        truth[correct, 0],
        truth[correct, 1],
        s=5.0,
        color=BLUE,
        edgecolor=WHITE,
        linewidth=0.15,
        alpha=0.82,
        zorder=4,
    )
    ax.scatter(
        truth[wrong, 0],
        truth[wrong, 1],
        s=7.0,
        color=RED,
        marker="x",
        linewidth=0.65,
        alpha=0.86,
        zorder=5,
    )
    ax.set_xlim(bounds[0], bounds[1])
    ax.set_ylim(bounds[2], bounds[3])
    ax.set_aspect("equal", adjustable="box")
    ax.set_xticks([])
    ax.set_yticks([])
    for spine in ax.spines.values():
        spine.set_visible(False)
    recall = 100.0 * correct.mean()
    ax.text(
        0.03,
        0.97,
        f"R@1 = {recall:.1f}%",
        transform=ax.transAxes,
        ha="left",
        va="top",
        fontsize=7.1,
        weight="bold",
        color=INK,
        bbox={"facecolor": WHITE, "edgecolor": "none", "alpha": 0.86, "pad": 1.2},
    )
    legend = [
        Line2D([], [], marker="o", ls="", color=BLUE, markersize=3.5, label="correct query"),
        Line2D([], [], marker="x", ls="", color=RED, markersize=3.5, label="wrong Top-1"),
    ]
    ax.legend(
        handles=legend,
        loc="lower left",
        bbox_to_anchor=(0.00, 0.00),
        ncol=1,
        borderaxespad=0.2,
        handletextpad=0.35,
        labelspacing=0.25,
    )
    save(fig, stem)


def make_rank_panel(loaded: dict[str, dict[str, list[dict[str, str]]]]) -> dict[str, dict[str, np.ndarray]]:
    curves: dict[str, dict[str, np.ndarray]] = {}
    fig, axes = plt.subplots(2, 1, figsize=(2.52, 2.38), sharex=True, facecolor=WHITE)
    for ax, dataset in zip(axes, ("IH", "CH")):
        curves[dataset] = {}
        for method in METHODS:
            curve = recall_curve(loaded[dataset][method])
            curves[dataset][method] = curve
            ax.plot(
                np.arange(1, 6),
                curve,
                color=METHODS[method],
                lw=1.75 if method == "UpDown-SC" else 1.0,
                marker="o",
                markersize=2.8 if method == "UpDown-SC" else 2.0,
                zorder=4 if method == "UpDown-SC" else 2,
                label=method,
            )
        ax.set_xlim(1, 5)
        ax.set_ylim(35 if dataset == "CH" else 50, 95)
        ax.set_ylabel("Recall (%)")
        ax.grid(axis="y", color=LIGHT, lw=0.55)
        ax.text(
            0.02,
            0.91,
            f"{dataset}  ($n_q={DATASETS[dataset]['n']}$)",
            transform=ax.transAxes,
            ha="left",
            va="top",
            fontsize=6.7,
            weight="bold",
            color=INK,
        )
    axes[-1].set_xlabel("retrieval rank $K$")
    axes[-1].set_xticks(np.arange(1, 6))
    handles, labels = axes[0].get_legend_handles_labels()
    fig.legend(
        handles,
        labels,
        loc="upper center",
        bbox_to_anchor=(0.54, 1.005),
        ncol=3,
        columnspacing=0.7,
        handlelength=1.4,
        handletextpad=0.3,
        borderaxespad=0.0,
    )
    fig.subplots_adjust(left=0.19, right=0.99, bottom=0.17, top=0.84, hspace=0.24)
    save(fig, "retrieval_rank_curves")
    return curves


def write_source_data(
    loaded: dict[str, dict[str, list[dict[str, str]]]],
    curves: dict[str, dict[str, np.ndarray]],
) -> None:
    with (OUT / "retrieval_distribution_source_data.csv").open("w", newline="") as handle:
        writer = csv.writer(handle)
        writer.writerow(
            [
                "record_type",
                "dataset",
                "method",
                "query_id_or_k",
                "truth_x",
                "truth_y",
                "top1_x",
                "top1_y",
                "first_correct_rank",
                "recall_percent",
            ]
        )
        for method in ("SC", "UpDown-SC"):
            for row in loaded["IH"][method]:
                writer.writerow(
                    [
                        "spatial_top1",
                        "IH",
                        method,
                        row["query_id"],
                        row["truth_x"],
                        row["truth_y"],
                        row["top1_x"],
                        row["top1_y"],
                        rank_value(row),
                        "",
                    ]
                )
        for dataset, dataset_curves in curves.items():
            for method, curve in dataset_curves.items():
                for k, value in enumerate(curve, start=1):
                    writer.writerow(
                        ["recall_curve", dataset, method, k, "", "", "", "", "", f"{value:.6f}"]
                    )

    metadata = {
        "overview": "Conceptual method overview; it contains no performance data.",
        "spatial_panels": {
            "dataset": "IH pilot",
            "queries": 320,
            "selection": "all eligible 2 m query keyframes",
            "encoding": "blue circle=correct Top-1; red cross/line=incorrect Top-1",
        },
        "rank_panel": {
            "datasets": {"IH": 320, "CH": 148},
            "metric": "Recall@K for K=1,...,5",
            "missing_or_empty_candidates": "counted as failures",
            "front_end": "gravity-canonicalized (+G) for all plotted methods",
        },
    }
    (OUT / "scancontext_style_figures_metadata.json").write_text(
        json.dumps(metadata, indent=2) + "\n"
    )


def add_text(slide, x, y, w, h, text, size=11, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_line(slide, x0, y0, x1, y1, color, width=1.0, transparency=0):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x0), Inches(y0), Inches(x1), Inches(y1)
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    line.line.transparency = transparency
    return line


def add_rect(slide, x, y, w, h, fill, line=LIGHT, width=0.8):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(width)
    return shape


def add_circle(slide, x, y, diameter, fill, line=WHITE, width=0.4):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x - diameter / 2),
        Inches(y - diameter / 2),
        Inches(diameter),
        Inches(diameter),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(width)
    return shape


def make_ppt(loaded: dict[str, dict[str, list[dict[str, str]]]], curves) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1 recreates the page-one schematic entirely with native shapes.
    slide = prs.slides.add_slide(blank)
    add_text(slide, 0.55, 0.16, 12.2, 0.45, "Measured indoor scan (side view)", 22, INK, True, PP_ALIGN.CENTER)
    slide.shapes.add_picture(
        str(OUT / "updown_sc_hero_sideview.png"),
        Inches(3.62),
        Inches(0.66),
        width=Inches(6.10),
        height=Inches(2.05),
    )

    add_rect(slide, 0.55, 3.10, 5.95, 3.85, "#FFF8F7", "#E7C6C3", 1.2)
    add_rect(slide, 6.83, 3.10, 5.95, 3.85, "#F5F9FD", "#BED4E8", 1.2)
    add_text(slide, 0.75, 3.20, 5.55, 0.42, "Conventional SC", 19, INK, True, PP_ALIGN.CENTER)
    add_text(slide, 7.03, 3.20, 5.55, 0.42, "UpDown-SC", 19, INK, True, PP_ALIGN.CENTER)

    for x in np.linspace(1.25, 5.75, 7):
        add_rect(slide, float(x), 3.95, 0.34, 1.38, LIGHT, LIGHT, 0.2)
        add_rect(slide, float(x), 3.95, 0.34, 0.46, UPPER, UPPER, 0.2)
    add_line(slide, 1.05, 4.10, 6.05, 4.10, UPPER, 2.0)
    add_text(slide, 1.05, 5.22, 5.0, 0.30, "max height", 12, UPPER, False, PP_ALIGN.CENTER)
    for i in range(10):
        add_rect(slide, 1.12 + i * 0.49, 5.68, 0.47, 0.55, UPPER, WHITE, 0.2)
    add_text(slide, 1.05, 6.35, 5.0, 0.30, "ceiling-dominated", 13, RED, True, PP_ALIGN.CENTER)

    split_y = 4.77
    add_line(slide, 7.43, split_y, 12.23, split_y, ORANGE, 1.5)
    add_text(slide, 12.25, split_y - 0.15, 0.45, 0.30, "τg", 12, ORANGE)
    heights = [0.75, 1.02, 0.63, 1.15, 0.83, 1.05, 0.68]
    for x, height in zip(np.linspace(7.70, 11.90, 7), heights):
        add_rect(slide, float(x), 5.45 - height, 0.34, height, "#4D8ABE", "#4D8ABE", 0.2)
        add_rect(slide, float(x), 3.95, 0.34, 0.45, UPPER, UPPER, 0.2)
    patterns = ([0, 1, 3, 2, 0, 1, 2, 3, 1, 0], [1, 2, 3, 1, 0, 2, 3, 2, 1, 0])
    for row, pattern in enumerate(patterns):
        # PowerPoint y grows downward, so row 0 is the visually upper row.
        color = UPPER if row == 0 else BLUE
        for i, level in enumerate(pattern):
            shade = mpl.colors.to_hex(
                np.array(mpl.colors.to_rgb(WHITE)) * (1 - 0.25 * level)
                + np.array(mpl.colors.to_rgb(color)) * (0.25 * level)
            )
            add_rect(slide, 7.43 + i * 0.48, 5.72 + row * 0.48, 0.46, 0.38, shade, WHITE, 0.2)
    add_text(slide, 7.43, 6.58, 4.80, 0.30, "lower structure retained", 13, BLUE, True, PP_ALIGN.CENTER)

    # Slide 2 recreates the aggregate spatial evidence with native shapes.
    slide = prs.slides.add_slide(blank)
    add_text(slide, 0.45, 0.15, 12.4, 0.45, "All-query retrieval evidence", 21, INK, True)
    map_xy = np.loadtxt(IH_MAP)[:, 1:3]
    truth_sc, top1_sc, correct_sc = spatial_arrays(loaded["IH"]["SC"])
    truth_ours, top1_ours, correct_ours = spatial_arrays(loaded["IH"]["UpDown-SC"])
    bounds = common_bounds(map_xy, np.vstack((truth_sc, truth_ours)), np.vstack((top1_sc, top1_ours)))

    def map_xy_to_slide(points, x, y, w, h):
        xmin, xmax, ymin, ymax = bounds
        px = x + (points[:, 0] - xmin) / (xmax - xmin) * w
        py = y + h - (points[:, 1] - ymin) / (ymax - ymin) * h
        return np.column_stack((px, py))

    for x, title, rows in (
        (0.45, "SC + gravity", loaded["IH"]["SC"]),
        (4.25, "UpDown-SC + gravity", loaded["IH"]["UpDown-SC"]),
    ):
        add_text(slide, x, 0.72, 3.45, 0.35, title, 14, INK, True, PP_ALIGN.CENTER)
        truth, top1, correct = spatial_arrays(rows)
        map_s = map_xy_to_slide(map_xy, x, 1.15, 3.45, 4.95)
        truth_s = map_xy_to_slide(truth, x, 1.15, 3.45, 4.95)
        top1_s = map_xy_to_slide(top1, x, 1.15, 3.45, 4.95)
        for p0, p1 in zip(map_s[:-1], map_s[1:]):
            add_line(slide, *p0, *p1, LIGHT, 0.45)
        for q, m, ok in zip(truth_s, top1_s, correct):
            if not ok:
                add_line(slide, *q, *m, RED, 0.3, 78)
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(q[0] - 0.018),
                Inches(q[1] - 0.018),
                Inches(0.036),
                Inches(0.036),
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = rgb(BLUE if ok else RED)
            dot.line.fill.background()
        add_text(
            slide,
            x + 0.10,
            5.98,
            1.55,
            0.32,
            f"R@1 = {100 * correct.mean():.1f}%",
            12,
            INK,
            True,
        )

    # Native-shape rank curves at right.
    x0, y0, w, h = 8.15, 1.15, 4.65, 2.15
    for dataset_index, dataset in enumerate(("IH", "CH")):
        yy = y0 + dataset_index * 2.75
        add_text(slide, x0, yy - 0.35, w, 0.30, f"{dataset} Recall@K", 13, INK, True)
        ymin = 50 if dataset == "IH" else 35
        ymax = 95
        for tick in (40, 60, 80) if dataset == "CH" else (60, 70, 80, 90):
            ty = yy + h - (tick - ymin) / (ymax - ymin) * h
            add_line(slide, x0, ty, x0 + w, ty, LIGHT, 0.45)
            add_text(slide, x0 - 0.42, ty - 0.12, 0.35, 0.23, str(tick), 8, MID, False, PP_ALIGN.RIGHT)
        for method, curve in curves[dataset].items():
            points = []
            for k, value in enumerate(curve, start=1):
                px = x0 + (k - 1) / 4 * w
                py = yy + h - (value - ymin) / (ymax - ymin) * h
                points.append((px, py))
            for p0, p1 in zip(points[:-1], points[1:]):
                add_line(slide, *p0, *p1, METHODS[method], 1.7 if method == "UpDown-SC" else 1.0)
        add_text(slide, x0, yy + h + 0.05, w, 0.24, "1          2          3          4          5", 8, MID, False, PP_ALIGN.CENTER)
    legend_y = 6.70
    for idx, method in enumerate(METHODS):
        lx = 8.15 + (idx % 3) * 1.50
        ly = legend_y + (idx // 3) * 0.28
        add_line(slide, lx, ly + 0.09, lx + 0.32, ly + 0.09, METHODS[method], 2.0 if method == "UpDown-SC" else 1.2)
        add_text(slide, lx + 0.37, ly, 1.05, 0.20, method, 8.2, INK)

    prs.save(OUT / "scancontext_style_figures_editable.pptx")


def main() -> None:
    loaded = load_all()
    hero_record, hero_xyz = load_hero_cloud()
    make_overview(hero_record, hero_xyz)
    map_xy = np.loadtxt(IH_MAP)[:, 1:3]
    truth_sc, top1_sc, _ = spatial_arrays(loaded["IH"]["SC"])
    truth_ours, top1_ours, _ = spatial_arrays(loaded["IH"]["UpDown-SC"])
    bounds = common_bounds(
        map_xy,
        np.vstack((truth_sc, truth_ours)),
        np.vstack((top1_sc, top1_ours)),
    )
    make_spatial_panel(map_xy, loaded["IH"]["SC"], bounds, "retrieval_spatial_sc")
    make_spatial_panel(
        map_xy, loaded["IH"]["UpDown-SC"], bounds, "retrieval_spatial_updown"
    )
    curves = make_rank_panel(loaded)
    write_source_data(loaded, curves)
    make_ppt(loaded, curves)
    print("Generated Scan Context-style paper figures and editable PowerPoint.")


if __name__ == "__main__":
    main()
