#!/usr/bin/env python3
"""Generate a measured-data retrieval example for the UpDown-SC paper.

Figure contract
---------------
Core conclusion:
    On an actual loop-corrected indoor map, the dual-envelope descriptor turns
    one deterministic SC Top-1 failure into a correct Top-1 retrieval.
Figure archetype:
    Image plate + spatial retrieval overlay.
Target/output:
    IEEE two-column paper; single-column PDF/SVG plus PNG preview and an
    editable-marker PowerPoint.
Backend:
    Python/Matplotlib.
Panel map:
    One full-map panel with map trajectory, query pseudo-reference, SC Top-1,
    UpDown-SC Top-1, and the predeclared 2 m positive region.
Evidence hierarchy:
    Hero evidence is the measured spatial error; the loop-corrected PCD is
    context and the mapping trajectory is the spatial reference.
Statistics:
    No inferential statistic. Distances are copied from per-query retrieval
    records and checked against the plotted coordinates.
Image-integrity notes:
    Every finite XY point inside the mapping-trajectory bounds contributes to
    the density raster. A global log1p transform and one global 98th-percentile
    display normalization are used; no spatial region is edited or sampled.
Reviewer risk:
    This is explicitly labelled a representative recovery case, not aggregate
    evidence; aggregate Recall@K remains in the accompanying table.
"""

from __future__ import annotations

import csv
import math
from pathlib import Path

import matplotlib as mpl
import matplotlib.pyplot as plt
import numpy as np
from matplotlib.lines import Line2D
from matplotlib.patches import Circle
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
EXPERIMENT = Path("${UPDOWN_SC_ROOT}/icra2027_runtime/experiments/private_two_bag_2m")
MAP_PCD = Path("${UPDOWN_SC_ROOT}/OneDrive/icra2027/slam/fast_lio/prior_map/scans.pcd")
MAP_POSES = EXPERIMENT / "map/session/optimized_poses_tum.txt"
SC_CSV = Path(
    "${UPDOWN_SC_ROOT}/icra2027_runtime/experiments/"
    "gravity_transfer_2m_20260718/ih/results/per_query.csv"
)
UPDOWN_CSV = Path(
    "${UPDOWN_SC_ROOT}/icra2027_runtime/experiments/"
    "updown_weight_ablation_real_20260721/selected/ih_gravity/paper_per_query.csv"
)
QUERY_ID = 151
POSITIVE_RADIUS_M = 2.0

BG = "#07111D"
MAP = "#AFC1CC"
TRAJECTORY = "#D9E2E8"
QUERY = "#2EC4B6"
SC = "#E58A22"
OURS = "#438BDE"
WHITE = "#F5F7F9"
MID = "#B8C5CF"


mpl.rcParams.update(
    {
        "font.family": "sans-serif",
        "font.sans-serif": ["Arial", "Helvetica", "DejaVu Sans", "sans-serif"],
        "svg.fonttype": "none",
        "pdf.fonttype": 42,
        "font.size": 7.5,
        "axes.linewidth": 0.8,
        "legend.frameon": False,
    }
)


def load_pcd_xy(path: Path) -> np.ndarray:
    with path.open("rb") as handle:
        header: dict[str, list[str]] = {}
        while True:
            line = handle.readline()
            if not line:
                raise RuntimeError(f"Missing DATA header in {path}")
            text = line.decode("ascii").strip()
            if not text or text.startswith("#"):
                continue
            key, *values = text.split()
            header[key.upper()] = values
            if key.upper() == "DATA":
                break
        if header["DATA"][0].lower() != "binary":
            raise RuntimeError(f"Only binary PCD is supported: {path}")
        fields = header["FIELDS"]
        sizes = [int(value) for value in header["SIZE"]]
        types = header["TYPE"]
        counts = [int(value) for value in header.get("COUNT", ["1"] * len(fields))]
        if any(count != 1 for count in counts):
            raise RuntimeError("Packed multi-count PCD fields are not supported")
        dtype_fields = []
        type_map = {
            ("F", 4): "<f4",
            ("F", 8): "<f8",
            ("I", 1): "<i1",
            ("I", 2): "<i2",
            ("I", 4): "<i4",
            ("U", 1): "<u1",
            ("U", 2): "<u2",
            ("U", 4): "<u4",
        }
        for field, size, kind in zip(fields, sizes, types):
            dtype_fields.append((field, type_map[(kind.upper(), size)]))
        points = np.fromfile(
            handle,
            dtype=np.dtype(dtype_fields),
            count=int(header["POINTS"][0]),
        )
    xy = np.column_stack((points["x"], points["y"])).astype(float, copy=False)
    return xy[np.isfinite(xy).all(axis=1)]


def load_retrieval(path: Path, algorithm: str) -> dict[str, float]:
    with path.open(newline="") as handle:
        for row in csv.DictReader(handle):
            if row.get("algorithm") == algorithm and int(row["query_id"]) == QUERY_ID:
                return {
                    key: float(row[key])
                    for key in (
                        "truth_x",
                        "truth_y",
                        "top1_x",
                        "top1_y",
                        "top1_error_m",
                    )
                }
    raise RuntimeError(f"Query {QUERY_ID} / {algorithm} is absent from {path}")


def checked_inputs():
    map_xy = np.loadtxt(MAP_POSES)[:, 1:3]
    pcd_xy = load_pcd_xy(MAP_PCD)
    sc = load_retrieval(SC_CSV, "SC + G")
    ours = load_retrieval(UPDOWN_CSV, "UpDown-SC pure adaptive")
    truth = np.array([sc["truth_x"], sc["truth_y"]])
    if not np.allclose(truth, [ours["truth_x"], ours["truth_y"]], atol=1e-6):
        raise RuntimeError("SC and UpDown-SC use different query references")
    for record in (sc, ours):
        top1 = np.array([record["top1_x"], record["top1_y"]])
        measured = float(np.linalg.norm(top1 - truth))
        if not math.isclose(measured, record["top1_error_m"], abs_tol=2e-3):
            raise RuntimeError("Recorded Top-1 error disagrees with plotted coordinates")
    return map_xy, pcd_xy, truth, sc, ours


def map_bounds(map_xy: np.ndarray) -> tuple[float, float, float, float]:
    margin = 3.0
    minimum = map_xy.min(axis=0) - margin
    maximum = map_xy.max(axis=0) + margin
    return minimum[0], maximum[0], minimum[1], maximum[1]


def density_image(
    pcd_xy: np.ndarray,
    bounds: tuple[float, float, float, float],
    bins: tuple[int, int] = (1100, 850),
) -> tuple[np.ndarray, np.ndarray]:
    xmin, xmax, ymin, ymax = bounds
    keep = (
        (pcd_xy[:, 0] >= xmin)
        & (pcd_xy[:, 0] <= xmax)
        & (pcd_xy[:, 1] >= ymin)
        & (pcd_xy[:, 1] <= ymax)
    )
    density, _, _ = np.histogram2d(
        pcd_xy[keep, 0],
        pcd_xy[keep, 1],
        bins=bins,
        range=((xmin, xmax), (ymin, ymax)),
    )
    density = np.log1p(density.T)
    positive = density[density > 0]
    scale = np.percentile(positive, 98) if positive.size else 1.0
    intensity = np.clip(density / max(scale, 1e-9), 0.0, 1.0) ** 0.48
    bg = np.array(mpl.colors.to_rgb(BG))
    fg = np.array(mpl.colors.to_rgb(MAP))
    rgb = bg + intensity[..., None] * (fg - bg)
    return rgb, keep


def make_figure() -> None:
    map_xy, pcd_xy, truth, sc, ours = checked_inputs()
    bounds = map_bounds(map_xy)
    raster, keep = density_image(pcd_xy, bounds)
    xmin, xmax, ymin, ymax = bounds
    sc_xy = np.array([sc["top1_x"], sc["top1_y"]])
    ours_xy = np.array([ours["top1_x"], ours["top1_y"]])

    fig, ax = plt.subplots(figsize=(3.45, 3.05), facecolor="white")
    ax.set_facecolor(BG)
    ax.imshow(
        raster,
        extent=(xmin, xmax, ymin, ymax),
        origin="lower",
        interpolation="bilinear",
        aspect="equal",
    )
    ax.plot(
        map_xy[:, 0],
        map_xy[:, 1],
        color=TRAJECTORY,
        linewidth=0.75,
        alpha=0.78,
        zorder=3,
    )
    ax.add_patch(
        Circle(
            truth,
            POSITIVE_RADIUS_M,
            facecolor="none",
            edgecolor=QUERY,
            linewidth=1.0,
            linestyle=(0, (3, 2)),
            zorder=5,
        )
    )
    ax.plot(
        [truth[0], sc_xy[0]],
        [truth[1], sc_xy[1]],
        color=SC,
        linewidth=1.35,
        alpha=0.95,
        zorder=5,
    )
    ax.plot(
        [truth[0], ours_xy[0]],
        [truth[1], ours_xy[1]],
        color=OURS,
        linewidth=1.55,
        alpha=0.95,
        zorder=6,
    )
    ax.scatter(
        *truth,
        s=38,
        marker="*",
        facecolor=QUERY,
        edgecolor=WHITE,
        linewidth=0.55,
        zorder=8,
    )
    ax.scatter(
        *sc_xy,
        s=35,
        marker="x",
        color=SC,
        linewidth=1.5,
        zorder=8,
    )
    ax.scatter(
        *ours_xy,
        s=26,
        marker="o",
        facecolor=OURS,
        edgecolor=WHITE,
        linewidth=0.55,
        zorder=8,
    )

    legend = [
        Line2D([0], [0], color=TRAJECTORY, lw=1.2, label="mapping trajectory"),
        Line2D(
            [0],
            [0],
            marker="*",
            color="none",
            markerfacecolor=QUERY,
            markeredgecolor=WHITE,
            markersize=7,
            label="query pseudo-reference",
        ),
        Line2D(
            [0],
            [0],
            marker="x",
            color=SC,
            lw=0,
            markersize=6,
            label=f"SC Top-1 ({sc['top1_error_m']:.1f} m)",
        ),
        Line2D(
            [0],
            [0],
            marker="o",
            color="none",
            markerfacecolor=OURS,
            markeredgecolor=WHITE,
            markersize=5.5,
            label=f"UpDown-SC Top-1 ({ours['top1_error_m']:.2f} m)",
        ),
    ]
    leg = ax.legend(
        handles=legend,
        loc="upper right",
        fontsize=6.4,
        labelcolor=WHITE,
        handlelength=1.4,
        borderpad=0.45,
        handletextpad=0.55,
        frameon=True,
        facecolor="#102033",
        edgecolor="#526579",
        framealpha=0.94,
    )
    leg.get_frame().set_linewidth(0.5)

    scale_length = 10.0
    scale_x0 = xmin + 3.0
    scale_y = ymin + 2.2
    ax.plot(
        [scale_x0, scale_x0 + scale_length],
        [scale_y, scale_y],
        color=WHITE,
        linewidth=1.8,
        solid_capstyle="butt",
        zorder=10,
    )
    ax.text(
        scale_x0 + scale_length / 2,
        scale_y + 0.75,
        "10 m",
        color=WHITE,
        fontsize=6.5,
        ha="center",
        va="bottom",
        zorder=10,
    )
    ax.set_xlim(xmin, xmax)
    ax.set_ylim(ymin, ymax)
    ax.set_xticks([])
    ax.set_yticks([])
    for spine in ax.spines.values():
        spine.set_visible(False)
    fig.subplots_adjust(left=0.01, right=0.99, bottom=0.01, top=0.99)

    export_options = {
        "facecolor": "white",
        "bbox_inches": "tight",
        "pad_inches": 0.015,
    }
    fig.savefig(OUT_DIR / "real_retrieval_example.svg", **export_options)
    fig.savefig(OUT_DIR / "real_retrieval_example.pdf", **export_options)
    fig.savefig(
        OUT_DIR / "real_retrieval_example.png", dpi=600, **export_options
    )
    plt.close(fig)

    print(
        "Wrote real retrieval figure: "
        f"query={QUERY_ID}, points={len(pcd_xy)}, displayed={int(keep.sum())}, "
        f"SC={sc['top1_error_m']:.3f} m, UpDown-SC={ours['top1_error_m']:.3f} m"
    )


def rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[:2], 16), int(value[2:4], 16), int(value[4:], 16))


def add_ppt_text(slide, x, y, w, h, text, size, color=WHITE, bold=False):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.margin_left = frame.margin_right = Inches(0.01)
    frame.margin_top = frame.margin_bottom = Inches(0)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shape


def make_editable_ppt() -> None:
    """Create a PPT with a raster PCD plate and editable scientific overlays."""
    map_xy, pcd_xy, truth, sc, ours = checked_inputs()
    bounds = map_bounds(map_xy)
    xmin, xmax, ymin, ymax = bounds
    raster, _ = density_image(pcd_xy, bounds)
    background_path = OUT_DIR / "real_retrieval_map_background.png"

    # Render only the dense measured map and mapping trajectory as the raster
    # plate. Every explanatory overlay is added as a native PowerPoint object.
    slide_width, slide_height = 6.90, 6.10
    ppt_raster_scale = 3.0
    bg_fig = plt.figure(
        figsize=(
            slide_width * ppt_raster_scale,
            slide_height * ppt_raster_scale,
        ),
        facecolor=BG,
    )
    bg_ax = bg_fig.add_axes([0.0, 0.0, 1.0, 1.0])
    bg_ax.imshow(
        raster,
        extent=(xmin, xmax, ymin, ymax),
        origin="lower",
        interpolation="bilinear",
        aspect="equal",
    )
    bg_ax.plot(
        map_xy[:, 0],
        map_xy[:, 1],
        color=TRAJECTORY,
        linewidth=1.5,
        alpha=0.78,
    )
    bg_ax.set_xlim(xmin, xmax)
    bg_ax.set_ylim(ymin, ymax)
    bg_ax.axis("off")
    bg_fig.canvas.draw()

    def world_to_slide(point: np.ndarray) -> tuple[float, float]:
        px, py = bg_ax.transData.transform(point)
        return (
            px / bg_fig.dpi / ppt_raster_scale,
            slide_height - py / bg_fig.dpi / ppt_raster_scale,
        )

    bg_fig.savefig(
        background_path,
        facecolor=BG,
        bbox_inches=None,
        pad_inches=0,
    )
    plt.close(bg_fig)

    prs = Presentation()
    prs.slide_width = Inches(slide_width)
    prs.slide_height = Inches(slide_height)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("#FFFFFF")
    slide.shapes.add_picture(
        str(background_path),
        Inches(0),
        Inches(0),
        width=Inches(slide_width),
        height=Inches(slide_height),
    )

    def add_line(
        start: np.ndarray,
        end: np.ndarray,
        color: str,
        width: float,
    ):
        x0, y0 = world_to_slide(start)
        x1, y1 = world_to_slide(end)
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x0),
            Inches(y0),
            Inches(x1),
            Inches(y1),
        )
        line.line.color.rgb = rgb(color)
        line.line.width = Pt(width)
        return line

    sc_xy = np.array([sc["top1_x"], sc["top1_y"]])
    ours_xy = np.array([ours["top1_x"], ours["top1_y"]])
    add_line(truth, sc_xy, SC, 2.2)
    add_line(truth, ours_xy, OURS, 2.6)

    tx, ty = world_to_slide(truth)
    radius_x, _ = world_to_slide(
        np.array([truth[0] + POSITIVE_RADIUS_M, truth[1]])
    )
    radius_in = abs(radius_x - tx)
    support = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(tx - radius_in),
        Inches(ty - radius_in),
        Inches(2 * radius_in),
        Inches(2 * radius_in),
    )
    support.fill.background()
    support.line.color.rgb = rgb(QUERY)
    support.line.width = Pt(1.7)
    support.line.dash_style = 2

    star_size = 0.16
    star = slide.shapes.add_shape(
        MSO_SHAPE.STAR_5_POINT,
        Inches(tx - star_size / 2),
        Inches(ty - star_size / 2),
        Inches(star_size),
        Inches(star_size),
    )
    star.fill.solid()
    star.fill.fore_color.rgb = rgb(QUERY)
    star.line.color.rgb = rgb(WHITE)
    star.line.width = Pt(0.8)

    sx, sy = world_to_slide(sc_xy)
    cross_size = 0.16
    for start, end in (
        ((sx - cross_size / 2, sy - cross_size / 2), (sx + cross_size / 2, sy + cross_size / 2)),
        ((sx - cross_size / 2, sy + cross_size / 2), (sx + cross_size / 2, sy - cross_size / 2)),
    ):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(start[0]),
            Inches(start[1]),
            Inches(end[0]),
            Inches(end[1]),
        )
        line.line.color.rgb = rgb(SC)
        line.line.width = Pt(2.2)

    ox, oy = world_to_slide(ours_xy)
    ours_marker = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(ox - 0.065),
        Inches(oy - 0.065),
        Inches(0.13),
        Inches(0.13),
    )
    ours_marker.fill.solid()
    ours_marker.fill.fore_color.rgb = rgb(OURS)
    ours_marker.line.color.rgb = rgb(WHITE)
    ours_marker.line.width = Pt(0.8)

    legend_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.86),
        Inches(0.12),
        Inches(2.92),
        Inches(1.12),
    )
    legend_box.fill.solid()
    legend_box.fill.fore_color.rgb = rgb("#102033")
    legend_box.fill.transparency = 5
    legend_box.line.color.rgb = rgb("#526579")
    legend_box.line.width = Pt(0.7)
    legend_rows = (
        ("mapping trajectory", TRAJECTORY),
        ("query pseudo-reference", QUERY),
        (f"SC Top-1 ({sc['top1_error_m']:.1f} m)", SC),
        (f"UpDown-SC Top-1 ({ours['top1_error_m']:.2f} m)", OURS),
    )
    for index, (text, color) in enumerate(legend_rows):
        y = 0.23 + 0.245 * index
        swatch = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(4.02),
            Inches(y + 0.015),
            Inches(0.075),
            Inches(0.075),
        )
        swatch.fill.solid()
        swatch.fill.fore_color.rgb = rgb(color)
        swatch.line.color.rgb = rgb(WHITE)
        swatch.line.width = Pt(0.4)
        add_ppt_text(slide, 4.16, y - 0.015, 2.46, 0.14, text, 6.5)

    scale_start = np.array([xmin + 3.0, ymin + 2.2])
    scale_end = np.array([xmin + 13.0, ymin + 2.2])
    add_line(scale_start, scale_end, WHITE, 2.4)
    scale_center_x, scale_center_y = world_to_slide(
        np.array([xmin + 8.0, ymin + 3.5])
    )
    add_ppt_text(
        slide,
        scale_center_x - 0.3,
        scale_center_y - 0.08,
        0.6,
        0.16,
        "10 m",
        6.5,
        WHITE,
    )

    note = slide.notes_slide.notes_text_frame
    note.text = (
        "Measured-data figure. Only the dense PCD map and mapping trajectory "
        "are rasterized; match lines, markers, legend, and scale are native "
        "editable PowerPoint objects. Source query 151; "
        f"SC error {sc['top1_error_m']:.3f} m; "
        f"UpDown-SC error {ours['top1_error_m']:.3f} m; "
        f"bounds=({xmin:.3f},{xmax:.3f},{ymin:.3f},{ymax:.3f}); "
        f"query=({truth[0]:.3f},{truth[1]:.3f}); "
        f"map keyframes={len(map_xy)}."
    )
    prs.save(OUT_DIR / "real_retrieval_example_editable.pptx")


if __name__ == "__main__":
    make_figure()
    make_editable_ppt()
